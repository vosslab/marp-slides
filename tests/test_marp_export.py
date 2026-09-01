"""Focused tests for direct native Marp presentation export."""

# Standard Library
import pathlib
from unittest import mock

# PIP3 modules
import PIL.Image
import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Local Modules
from marp_lib import layouts
import marp_lib.native_export


HEADER = "---\nmarp: true\ntheme: genetics\nsize: 16:10\n---\n"


#============================================
def write_png(output_path: pathlib.Path) -> pathlib.Path:
	"""Write one small component image fixture."""
	image = PIL.Image.new("RGB", (80, 40), (36, 87, 143))
	image.save(output_path)
	return output_path


#============================================
def cells(count: int) -> str:
	"""Return valid reading-order quoted cells for a multi-content layout."""
	return "\n\n".join(f"> ## Cell {index}\n>\n> - Editable item {index}" for index in range(1, count + 1))


#============================================
def layout_markdown(name: str, spec: layouts.LayoutSpec) -> str:
	"""Return the smallest valid source fixture for one registered layout."""
	class_directive = f"<!-- _class: {name} -->\n"
	if name == "blank":
		return HEADER + class_directive
	if name in ("title-slide", "centered-text"):
		return HEADER + class_directive + "# Center title\n\n## Subtitle\n"
	if name == "title-only":
		return HEADER + class_directive + "# Title only\n"
	if name == "gallery":
		return HEADER + class_directive + "![One](one.png) ![Two](two.png)\n"
	if spec.cell_count:
		return HEADER + class_directive + "# Layout title\n\n" + cells(spec.cell_count) + "\n"
	return HEADER + class_directive + "# Layout title\n\n- Editable body\n"


#============================================
@pytest.mark.parametrize("name", sorted(layouts.LAYOUTS))
def test_every_registered_layout_renders_native_objects(tmp_path: pathlib.Path, name: str) -> None:
	"""Each canonical layout selects its own builder and retains editable objects."""
	spec = layouts.LAYOUTS[name]
	if name == "gallery":
		write_png(tmp_path / "one.png")
		write_png(tmp_path / "two.png")
	deck_path = tmp_path / f"{name}.md"
	deck_path.write_text(layout_markdown(name, spec), encoding="utf-8")
	output_path = tmp_path / f"{name}.pptx"
	marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path), output_path)
	slide = Presentation(output_path).slides[0]
	assert layouts.LAYOUTS[name] is spec
	assert all(shape.shape_type != MSO_SHAPE_TYPE.PICTURE or
		(shape.width < Presentation(output_path).slide_width and shape.height < Presentation(output_path).slide_height)
		for shape in slide.shapes)
	if name not in ("blank", "gallery"):
		assert any(shape.has_text_frame for shape in slide.shapes)
	if spec.cell_count:
		text = "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)
		assert f"Cell {spec.cell_count}" in text


#============================================
def test_horizontal_grid_layout_keeps_bullets_numbers_and_links(tmp_path: pathlib.Path) -> None:
	"""Two-content cells retain independently editable native text semantics."""
	deck_path = tmp_path / "two-content.md"
	deck_path.write_text(HEADER + "<!-- _class: title-two-content -->\n# Overview\n\n"
		"> ## [Left](https://example.edu/left)\n>\n> - First\n>   - Nested\n> 1. Ordered\n\n"
		"> ## Right\n>\n> - Second\n", encoding="utf-8")
	output_path = tmp_path / "two-content.pptx"
	marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path), output_path)
	shape_xml = "".join(shape.element.xml for shape in Presentation(output_path).slides[0].shapes)
	assert "buChar" in shape_xml and "buAutoNum" in shape_xml and "hlinkClick" in shape_xml


#============================================
def test_inline_runs_keep_native_formatting_and_url_typography(tmp_path: pathlib.Path) -> None:
	"""Typed emphasis, code, breaks, and bare URLs remain formatted editable runs."""
	deck_path = tmp_path / "formatting.md"
	deck_path.write_text(HEADER + "<!-- _class: title-content -->\n# [Heading](https://example.edu)\n\n"
		"*Italic* and `code` [Course resource](https://example.edu/resource)  \nhttps://example.edu/path\n", encoding="utf-8")
	output_path = tmp_path / "formatting.pptx"
	marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path), output_path)
	runs = [run for shape in Presentation(output_path).slides[0].shapes if shape.has_text_frame
		for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]
	assert any(run.text == "Italic" and run.font.italic for run in runs)
	assert any(run.text == "code" and run.font.name == layouts.FONT_NAME for run in runs)
	assert any(run.text == "Course resource" and run.font.name == layouts.FONT_NAME and
		run.hyperlink.address == "https://example.edu/resource" for run in runs)
	assert any(run.text == "https://example.edu/path" and run.font.name == layouts.URL_FONT_NAME
		and run.hyperlink.address == "https://example.edu/path" for run in runs)


#============================================
def test_h1_size_modifier_writes_150_point_title_without_changing_subtitle_or_pagination(
		tmp_path: pathlib.Path) -> None:
	"""font-size-200 remains an editable 150pt H1 and leaves other runs normal."""
	deck_path = tmp_path / "display-title.md"
	deck_path.write_text(HEADER + "<!-- _class: font-size-200 centered-text -->\n# THE END\n\n"
		"## Normal subtitle\n", encoding="utf-8")
	output_path = tmp_path / "display-title.pptx"
	marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path), output_path)
	runs = [run for shape in Presentation(output_path).slides[0].shapes if shape.has_text_frame
		for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]
	assert next(run for run in runs if run.text == "THE END").font.size.pt == 150
	assert next(run for run in runs if run.text == "Normal subtitle").font.size.pt == layouts.css_px_to_pt(31)
	assert next(run for run in runs if run.text == "1").font.size.pt == layouts.css_px_to_pt(18)


#============================================
def test_h1_size_modifier_reports_native_title_capacity(tmp_path: pathlib.Path) -> None:
	"""A requested display title fails rather than receiving an implicit size reduction."""
	deck_path = tmp_path / "too-large.md"
	deck_path.write_text(HEADER + "<!-- _class: title-content font-size-200 -->\n# Too large\n\n"
		"- Editable body\n", encoding="utf-8")
	with pytest.raises(ValueError, match=r"too-large\.md:7:.*title-content H1 font-size-200 does not fit"):
		marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path),
			tmp_path / "too-large.pptx")


#============================================
@pytest.mark.parametrize("name", ["title-vertical-text", "vertical-title-vertical-text",
	"title-two-vertical-text-clipart"])
def test_vertical_layouts_write_editable_ooxml_text_direction(tmp_path: pathlib.Path, name: str) -> None:
	"""Vertical layouts use editable text-body direction rather than a raster image."""
	spec = layouts.LAYOUTS[name]
	deck_path = tmp_path / f"{name}.md"
	deck_path.write_text(layout_markdown(name, spec), encoding="utf-8")
	output_path = tmp_path / f"{name}.pptx"
	marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path), output_path)
	shape_xml = "".join(shape.element.xml for shape in Presentation(output_path).slides[0].shapes)
	assert 'vert="vert"' in shape_xml


#============================================
def test_gallery_images_keep_descriptions_and_are_not_full_slide(tmp_path: pathlib.Path) -> None:
	"""Gallery components retain their independent descriptions after PPTX export."""
	write_png(tmp_path / "one.png")
	write_png(tmp_path / "two.png")
	deck_path = tmp_path / "gallery.md"
	deck_path.write_text(HEADER + "<!-- _class: gallery -->\n# Components\n\n"
		"![First component](one.png) ![Second component](two.png)\n", encoding="utf-8")
	output_path = tmp_path / "gallery.pptx"
	marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path), output_path)
	presentation = Presentation(output_path)
	pictures = [shape for shape in presentation.slides[0].shapes
		if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
	assert [picture.element.nvPicPr.cNvPr.get("descr") for picture in pictures] == [
		"First component", "Second component",
	]
	assert all(picture.width < presentation.slide_width and picture.height < presentation.slide_height
		for picture in pictures)


#============================================
def test_rejects_retired_layout_classes_directly(tmp_path: pathlib.Path) -> None:
	"""Retired layout and modifier vocabulary has no compatibility contract."""
	deck_path = tmp_path / "retired.md"
	deck_path.write_text(HEADER + "<!-- _class: lead -->\n# Retired\n", encoding="utf-8")
	with pytest.raises(ValueError, match="unsupported Marp slide class"):
		marp_lib.native_export.parse_deck(deck_path)


#============================================
def test_rejects_wrong_cell_count_and_mixed_cell_content(tmp_path: pathlib.Path) -> None:
	"""Cell source contracts account for every source block before rendering."""
	image_path = write_png(tmp_path / "component.png")
	wrong_count = tmp_path / "wrong-count.md"
	wrong_count.write_text(HEADER + "<!-- _class: title-four-content -->\n# Four\n\n" + cells(3),
		encoding="utf-8")
	with pytest.raises(ValueError, match=r"wrong-count\.md:\d+:.*require exactly 4"):
		marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(wrong_count),
			tmp_path / "wrong-count.pptx")
	mixed = tmp_path / "mixed.md"
	mixed.write_text(HEADER + "<!-- _class: title-two-content -->\n# Mixed\n\n"
		"> ## First\n>\n> - Text\n>\n> ![Component](component.png)\n\n> ## Second\n>\n> - Text\n",
		encoding="utf-8")
	with pytest.raises(ValueError, match=r"mixed\.md:13:.*cannot combine text"):
		marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(mixed),
			tmp_path / "mixed.pptx")
	assert image_path.is_file()


#============================================
def test_missing_component_image_reports_its_authored_image_line(tmp_path: pathlib.Path) -> None:
	"""Image-resolution failures retain the precise component image location."""
	deck_path = tmp_path / "missing-image.md"
	deck_path.write_text(HEADER + "<!-- _class: title-content -->\n# Missing\n\n"
		"![Absent component](absent.png)\n", encoding="utf-8")
	with pytest.raises(ValueError, match=r"missing-image\.md:9:.*component image is missing"):
		marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path),
			tmp_path / "missing-image.pptx")


#============================================
@pytest.mark.parametrize("name", ["title-vertical-text", "vertical-title-vertical-text"])
@pytest.mark.parametrize("body", ["Vertical paragraph\n", "- Vertical list item\n"])
def test_vertical_root_body_layouts_accept_one_text_block(tmp_path: pathlib.Path, name: str,
		body: str) -> None:
	"""Each root-body vertical layout creates exactly one editable vertical body frame."""
	deck_path = tmp_path / f"{name}.md"
	deck_path.write_text(HEADER + f"<!-- _class: {name} -->\n# Vertical\n\n" + body, encoding="utf-8")
	output_path = tmp_path / f"{name}.pptx"
	marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path), output_path)
	vertical_frames = [shape for shape in Presentation(output_path).slides[0].shapes if shape.has_text_frame
		and 'vert="vert"' in shape.element.xml and shape.text != "Vertical"]
	assert len(vertical_frames) == 1


#============================================
@pytest.mark.parametrize("name", ["title-vertical-text", "vertical-title-vertical-text"])
def test_vertical_root_body_layouts_accept_one_component_image(tmp_path: pathlib.Path, name: str) -> None:
	"""A single component image is the legal native picture form of a vertical body."""
	write_png(tmp_path / "component.png")
	deck_path = tmp_path / f"{name}-image.md"
	deck_path.write_text(HEADER + f"<!-- _class: {name} -->\n# Vertical\n\n"
		"![Vertical component](component.png)\n", encoding="utf-8")
	output_path = tmp_path / f"{name}-image.pptx"
	marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path), output_path)
	pictures = [shape for shape in Presentation(output_path).slides[0].shapes
		if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
	assert len(pictures) == 1


#============================================
@pytest.mark.parametrize("name", ["title-vertical-text", "vertical-title-vertical-text"])
def test_vertical_root_body_layouts_reject_multiple_blocks_at_the_second_block(tmp_path: pathlib.Path,
		name: str) -> None:
	"""A second root body block receives its own actionable source location."""
	deck_path = tmp_path / f"{name}-invalid.md"
	deck_path.write_text(HEADER + f"<!-- _class: {name} -->\n# Vertical\n\nFirst paragraph\n\n- Second block\n",
		encoding="utf-8")
	with pytest.raises(ValueError, match=rf"{name}-invalid\.md:11:.*exactly one root body block"):
		marp_lib.native_export.render_native_pptx(marp_lib.native_export.parse_deck(deck_path),
			tmp_path / "invalid.pptx")


#============================================
def test_layout_registry_has_one_individual_builder_per_layout() -> None:
	"""Every declared layout owns a distinct named native builder boundary."""
	builders = [spec.builder for spec in layouts.LAYOUTS.values()]
	assert len(builders) == len(set(builders))


#============================================
def test_presentation_chain_converts_odp_to_pdf(tmp_path: pathlib.Path) -> None:
	"""PDF export uses ODP as its required LibreOffice predecessor."""
	deck_path = tmp_path / "chain.md"
	deck_path.write_text(HEADER + "<!-- _class: title-content -->\n# Chain\n\n- Editable body\n", encoding="utf-8")
	with mock.patch.object(marp_lib.native_export, "find_repo_root", return_value=tmp_path), \
		mock.patch.object(marp_lib.native_export, "convert_presentation") as convert:
		outputs = marp_lib.native_export.export_deck(str(deck_path), "pdf")
	assert list(outputs) == ["pptx", "odp", "pdf"]
	assert convert.call_args_list[0].args[0] == outputs["pptx"]
	assert convert.call_args_list[0].args[2] == "odp"
	assert convert.call_args_list[1].args[0] == outputs["odp"]
	assert convert.call_args_list[1].args[2] == "pdf"
