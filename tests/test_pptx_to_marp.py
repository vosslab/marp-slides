"""Behavioral tests for structured PPTX-to-Marp conversion."""

# Standard Library
import json
import zipfile
import pathlib

# PIP3 modules
import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

# local repo modules
from marp_lib import native_export
from tools import pptx_to_marp


#============================================
def write_png(output_path: pathlib.Path, color: tuple[int, int, int]) -> pathlib.Path:
	"""Write one bounded raster fixture."""
	image = Image.new("RGB", (40, 30), color)
	image.save(output_path)
	return output_path


#============================================
def add_title(slide: object, title: str) -> None:
	"""Set the title placeholder on a synthetic slide."""
	title_shape = slide.shapes.title
	assert title_shape is not None
	title_shape.text = title


#============================================
def write_split_pptx(output_path: pathlib.Path, image_path: pathlib.Path) -> pathlib.Path:
	"""Write one visible split slide and one hidden slide."""
	presentation = Presentation()
	visible = presentation.slides.add_slide(presentation.slide_layouts[5])
	add_title(visible, "Genetics overview")
	text_box = visible.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(5.2), Inches(3.0))
	text_box.text_frame.text = "Chromosomes carry genes"
	visible.shapes.add_picture(
		str(image_path),
		Inches(7.0),
		Inches(1.5),
		width=Inches(2.5),
	)
	visible.notes_slide.notes_text_frame.text = "Explain inheritance -- then pause"
	hidden = presentation.slides.add_slide(presentation.slide_layouts[5])
	add_title(hidden, "Hidden source slide")
	hidden.element.set("show", "0")
	presentation.save(output_path)
	return output_path


#============================================
def test_structured_split_conversion_preserves_count_notes_and_visibility(
	tmp_path: pathlib.Path,
) -> None:
	"""Text, image, note, order, and hidden state survive without a slide screenshot."""
	image_path = write_png(tmp_path / "chromosome.png", (20, 90, 160))
	input_path = write_split_pptx(tmp_path / "lecture.pptx", image_path)
	output_path = tmp_path / "lecture.md"

	summary = pptx_to_marp.convert_pptx(
		input_path,
		output_path,
		expected_slide_count=2,
		expected_hidden={2},
	)
	markdown = output_path.read_text(encoding="utf-8")
	report = json.loads(summary.report_path.read_text(encoding="utf-8"))

	assert summary.visible_slides == 1
	assert summary.editable_slides == 1
	assert summary.hidden_slides == 1
	assert summary.extracted_images == 1
	assert "# Genetics overview" in markdown
	assert "Chromosomes carry genes" in markdown
	assert "<!-- _class: title-two-content -->" in markdown
	assert "> - Chromosomes carry genes" in markdown
	assert "> ![Slide image 1]" in markdown
	assert "bg right" not in markdown
	assert "Explain inheritance - - then pause" in markdown
	assert "source-fallback" not in markdown
	assert "slide_001_source.png" not in markdown
	assert report["hidden_slides"] == [2]
	assert len(report["slides"]) == 1
	assert report["slides"][0]["layout"] == "title-two-content"


#============================================
def test_three_images_use_one_auto_fitting_gallery_slide(tmp_path: pathlib.Path) -> None:
	"""Three source pictures remain one editable slide with a shared gallery layout."""
	presentation = Presentation()
	slide = presentation.slides.add_slide(presentation.slide_layouts[5])
	add_title(slide, "Microscope views")
	for index, color in enumerate(((100, 0, 0), (0, 100, 0), (0, 0, 100))):
		image_path = write_png(tmp_path / f"view_{index}.png", color)
		slide.shapes.add_picture(
			str(image_path),
			Inches(0.8 + index * 3.0),
			Inches(1.7),
			width=Inches(2.4),
		)
	input_path = tmp_path / "gallery.pptx"
	presentation.save(input_path)
	output_path = tmp_path / "gallery.md"

	summary = pptx_to_marp.convert_pptx(input_path, output_path)
	markdown = output_path.read_text(encoding="utf-8")

	assert summary.visible_slides == 1
	assert summary.extracted_images == 3
	assert "<!-- _class: gallery -->" in markdown
	assert markdown.count("![Slide image") == 3
	assert markdown.count("\n---\n") == 1


#============================================
def make_slide_data(
	*,
	source_index: int = 1,
	title: str = "Slide title",
	text_lines: tuple[tuple[int, str], ...] = (),
	image_positions: tuple[tuple[int, int], ...] = (),
) -> pptx_to_marp.SlideData:
	"""Build one extracted slide model for layout-classifier branch tests."""
	text_blocks = () if not text_lines else (
		pptx_to_marp.TextBlock(100, 200, text_lines, False),
	)
	images = tuple(
		pptx_to_marp.ImageAsset(left, top, 200, 120, f"assets/deck/image_{index}.png", "Image")
		for index, (left, top) in enumerate(image_positions, start=1)
	)
	return pptx_to_marp.SlideData(
		source_index, False, (title,), text_blocks, images, (), (),
	)


#============================================
def test_layout_classifier_emits_only_canonical_explicit_classes() -> None:
	"""Each supported importer source shape has one native layout declaration."""
	cases = (
		(
			"opening title", make_slide_data(), True, "title-slide",
			("<!-- _class: title-slide -->",),
		),
		(
			"later title only", make_slide_data(), False, "title-only",
			("<!-- _class: title-only -->",),
		),
		(
			"text body", make_slide_data(text_lines=((0, "Editable body"),)), False,
			"title-content", ("<!-- _class: title-content -->", "- Editable body"),
		),
		(
			"one image body", make_slide_data(image_positions=((700, 200),)), False,
			"title-content", ("<!-- _class: title-content -->", "![Image]"),
		),
		(
			"right image cell", make_slide_data(
				text_lines=((0, "Text first"),), image_positions=((800, 200),),
			), False, "title-two-content", ("> - Text first", "> ![Image]"),
		),
		(
			"left image cell", make_slide_data(
				text_lines=((0, "Text second"),), image_positions=((20, 200),),
			), False, "title-two-content", ("> ![Image]", "> - Text second"),
		),
		(
			"gallery", make_slide_data(image_positions=((100, 200), (400, 200), (700, 200))),
			False, "gallery", ("<!-- _class: gallery -->",),
		),
		(
			"multi image cell", make_slide_data(
				text_lines=((0, "Text cell"),), image_positions=((300, 200), (700, 200)),
			), False, "title-two-content", ("> - Text cell", "> ![Image]"),
		),
	)
	for _name, slide, is_first, expected_layout, expected_fragments in cases:
		lines, layout = pptx_to_marp.render_slide(slide, 1000, is_first)
		markdown = "\n".join(lines)
		assert layout == expected_layout
		assert markdown.count("<!-- _class:") == 1
		assert f"<!-- _class: {expected_layout} -->" in markdown
		for fragment in expected_fragments:
			assert fragment in markdown
		assert "lead" not in markdown
		assert "figure" not in markdown
		assert "two-pane" not in markdown
		assert "bg left" not in markdown
		assert "bg right" not in markdown
		if _name == "multi image cell":
			assert markdown.count("![Image]") == 2


#============================================
def test_imported_text_image_slide_parses_and_renders_natively(tmp_path: pathlib.Path) -> None:
	"""Importer Markdown is accepted by the repository-owned native renderer."""
	image_path = write_png(tmp_path / "component.png", (20, 90, 160))
	input_path = write_split_pptx(tmp_path / "source.pptx", image_path)
	markdown_path = tmp_path / "deck.md"
	pptx_to_marp.convert_pptx(input_path, markdown_path)

	deck = native_export.parse_deck(markdown_path)
	pptx_path = native_export.render_native_pptx(deck, tmp_path / "native.pptx")

	assert len(deck.slides) == 1
	assert pptx_path.is_file()
	assert Presentation(pptx_path).slides[0].shapes


#============================================
def test_visibility_disagreement_fails_closed(tmp_path: pathlib.Path) -> None:
	"""The normalized PPTX cannot silently override authoritative ODP visibility."""
	presentation = Presentation()
	slide = presentation.slides.add_slide(presentation.slide_layouts[5])
	add_title(slide, "Visible")
	input_path = tmp_path / "visible.pptx"
	presentation.save(input_path)

	with pytest.raises(RuntimeError, match="visibility disagree"):
		pptx_to_marp.convert_pptx(
			input_path,
			tmp_path / "visible.md",
			expected_slide_count=1,
			expected_hidden={1},
		)


#============================================
def test_unsafe_pptx_archive_member_is_rejected(tmp_path: pathlib.Path) -> None:
	"""OOXML archive traversal fails before python-pptx parsing."""
	input_path = tmp_path / "unsafe.pptx"
	with zipfile.ZipFile(input_path, "w") as archive:
		archive.writestr("[Content_Types].xml", "<Types/>")
		archive.writestr("ppt/presentation.xml", "<presentation/>")
		archive.writestr("../outside.png", b"not an image")

	with pytest.raises(ValueError, match="unsafe archive member path"):
		pptx_to_marp.validate_pptx(input_path)


#============================================
def test_converter_refuses_existing_canonical_output(tmp_path: pathlib.Path) -> None:
	"""A second import cannot replace established Markdown or assets."""
	output_path = tmp_path / "lecture.md"
	output_path.write_text("existing\n", encoding="utf-8")
	with pytest.raises(FileExistsError, match="will not overwrite"):
		pptx_to_marp.validate_output_path(output_path)

	output_path.unlink()
	(tmp_path / "assets" / "lecture").mkdir(parents=True)
	with pytest.raises(FileExistsError, match="asset directory"):
		pptx_to_marp.validate_output_path(output_path)
