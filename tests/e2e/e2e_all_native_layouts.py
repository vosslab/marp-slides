#!/usr/bin/env python3
"""Exercise the complete native layout catalog through PPTX, ODP, and PDF."""

# Standard Library
import pathlib
import re
import shutil
import subprocess
import sys
import tempfile
import uuid
import zipfile

# PIP3 modules
import defusedxml.ElementTree
import PIL.Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


NAMESPACES = {
	"draw": "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0",
	"fo": "urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0",
	"office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
	"presentation": "urn:oasis:names:tc:opendocument:xmlns:presentation:1.0",
	"svg": "urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0",
	"style": "urn:oasis:names:tc:opendocument:xmlns:style:1.0",
	"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
	"xlink": "http://www.w3.org/1999/xlink",
}
LAYOUTS = (
	"blank", "title-only", "title-slide", "title-content", "centered-text",
	"title-two-content", "title-content-and-two-content",
	"title-two-content-and-content", "title-content-over-content",
	"title-two-content-over-content", "title-four-content", "title-six-content",
	"vertical-title-vertical-text", "vertical-title-text-chart", "title-vertical-text",
	"title-two-vertical-text-clipart", "gallery",
)
VERTICAL_LAYOUTS = frozenset((
	"vertical-title-vertical-text", "vertical-title-text-chart", "title-vertical-text",
	"title-two-vertical-text-clipart",
))
MAX_COMPONENT_PAGE_FRACTION = 0.98
ODF_LENGTH_TO_CENTIMETERS = {
	"cm": 1.0, "mm": 0.1, "in": 2.54, "pt": 2.54 / 72.0,
	"pc": 2.54 / 6.0, "px": 2.54 / 96.0,
}


#============================================
def repo_root() -> pathlib.Path:
	"""Return the repository root containing this acceptance runner."""
	return pathlib.Path(__file__).resolve().parents[2]


#============================================
def require(condition: bool, message: str) -> None:
	"""Raise an actionable failure when a native-export contract is absent."""
	if not condition:
		raise RuntimeError(message)


#============================================
def write_component_images(workspace: pathlib.Path) -> None:
	"""Write small source-owned component images with distinct solid colors."""
	for index, color in enumerate(((36, 87, 143), (92, 161, 70), (204, 116, 47)), start=1):
		image = PIL.Image.new("RGB", (160, 96), color)
		image.save(workspace / f"component_{index}.png")


#============================================
def cell(label: str) -> str:
	"""Return one blockquote cell with a heading, linked list, and native text."""
	return (
		f"> ## {label}\n>\n> - Editable {label} [link](https://example.edu/{label.lower()})\n"
		">   - Nested editable detail\n"
	)


#============================================
def image_cell(label: str, image_name: str) -> str:
	"""Return one blockquote cell with a descriptive native component image."""
	return f"> ## {label}\n>\n> ![Native {label} component image]({image_name})\n"


#============================================
def multi_cell_slide(layout: str, count: int) -> str:
	"""Return a valid grid-layout slide with reading-order blockquote cells."""
	cells = [cell(f"{layout} cell {index}") for index in range(1, count + 1)]
	if layout in ("title-two-content", "vertical-title-text-chart",
		"title-two-vertical-text-clipart"):
		cells[-1] = image_cell(f"{layout} image", "component_1.png")
	return f"<!-- _class: {layout} -->\n# {layout} acceptance\n\n" + "\n".join(cells)


#============================================
def write_deck(deck_path: pathlib.Path) -> None:
	"""Create one valid canonical slide for every supported native layout."""
	slides = [
		"<!-- _class: blank -->",
		"<!-- _class: title-only -->\n# title-only acceptance",
		"<!-- _class: title-slide -->\n# title-slide acceptance\n\n## Native object model",
		"<!-- _class: title-content -->\n# title-content acceptance\n\n"
		"- Editable [resource](https://example.edu/title-content)\n  - Nested editable detail\n"
		"1. Ordered editable step\n\n<!-- notes: Present this native title-content slide. -->",
		"<!-- _class: centered-text font-size-200 -->\n# THE END\n\n## Editable interstitial",
		multi_cell_slide("title-two-content", 2),
		multi_cell_slide("title-content-and-two-content", 3),
		multi_cell_slide("title-two-content-and-content", 3),
		multi_cell_slide("title-content-over-content", 2),
		multi_cell_slide("title-two-content-over-content", 3),
		multi_cell_slide("title-four-content", 4),
		multi_cell_slide("title-six-content", 6),
		"<!-- _class: vertical-title-vertical-text -->\n# vertical-title-vertical-text acceptance\n\n"
		"- Vertical editable [resource](https://example.edu/vertical-title-vertical-text)",
		multi_cell_slide("vertical-title-text-chart", 2),
		"<!-- _class: title-vertical-text -->\n# title-vertical-text acceptance\n\n"
		"- Vertical body [resource](https://example.edu/title-vertical-text)",
		multi_cell_slide("title-two-vertical-text-clipart", 3),
		"<!-- _class: gallery -->\n# gallery acceptance\n\n"
		"![Native gallery component one](component_1.png)\n"
		"![Native gallery component two](component_2.png)\n"
		"![Native gallery component three](component_3.png)",
	]
	require(len(slides) == len(LAYOUTS), "Acceptance deck contains exactly one slide per layout")
	deck_path.write_text(
		"---\nmarp: true\ntheme: genetics\nsize: 16:10\ntitle: Native layout catalog\n---\n" +
		"\n---\n".join(slides) + "\n",
		encoding="utf-8",
	)


#============================================
def slide_text(slide: object) -> str:
	"""Return editable text from all text frames on one PPTX slide."""
	return "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)


#============================================
def inspect_pptx(pptx_path: pathlib.Path) -> None:
	"""Verify native PPTX text, list, link, notes, picture, and vertical objects."""
	presentation = Presentation(pptx_path)
	require(len(presentation.slides) == len(LAYOUTS),
		f"PPTX has {len(presentation.slides)} slides; expected {len(LAYOUTS)} layouts")
	all_xml = "\n".join("".join(shape.element.xml for shape in slide.shapes)
		for slide in presentation.slides)
	all_text = "\n".join(slide_text(slide) for slide in presentation.slides)
	require("title-content acceptance" in all_text and "THE END" in all_text and "gallery acceptance" in all_text,
		"PPTX preserves editable text across the layout catalog")
	title_runs = [run for shape in presentation.slides[4].shapes if shape.has_text_frame
		for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text == "THE END"]
	require(len(title_runs) == 1 and title_runs[0].font.size.pt == 150,
		"PPTX preserves font-size-200 as an editable 150pt title")
	require("buChar" in all_xml and "buAutoNum" in all_xml,
		"PPTX preserves native unordered and ordered list structures")
	require("hlinkClick" in all_xml, "PPTX preserves native hyperlink relationships")
	require(presentation.slides[3].notes_slide.notes_text_frame.text ==
		"Present this native title-content slide.", "PPTX preserves presenter notes")
	for layout in VERTICAL_LAYOUTS:
		index = LAYOUTS.index(layout)
		slide_xml = "".join(shape.element.xml for shape in presentation.slides[index].shapes)
		require('vert="vert"' in slide_xml,
			f"PPTX {layout} uses OOXML vertical text direction")
	pictures = [shape for slide in presentation.slides for shape in slide.shapes
		if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
	require(len(pictures) >= 6, "PPTX retains catalog component images as picture objects")
	for picture in pictures:
		description = picture.element.nvPicPr.cNvPr.get("descr", "")
		require(description.startswith("Native "),
			"PPTX component pictures preserve authored image descriptions")
		require(picture.width < presentation.slide_width * MAX_COMPONENT_PAGE_FRACTION or
			picture.height < presentation.slide_height * MAX_COMPONENT_PAGE_FRACTION,
			"PPTX contains no full-slide picture fallback")


#============================================
def text_content(element: object | None) -> str:
	"""Return all descendant text nodes from one ODF XML element."""
	return "" if element is None else "".join(element.itertext())


#============================================
def parse_odf_length(length_text: str) -> float:
	"""Return a supported ODF physical length in centimeters."""
	match = re.fullmatch(r"([+-]?(?:\d+(?:\.\d*)?|\.\d+))(cm|mm|in|pt|pc|px)", length_text)
	require(match is not None, f"ODP length {length_text!r} has a supported physical unit")
	return float(match.group(1)) * ODF_LENGTH_TO_CENTIMETERS[match.group(2)]


#============================================
def page_dimensions(styles_root: object, page: object) -> tuple[float, float]:
	"""Return page dimensions in centimeters for the ODP page's master layout."""
	style_name = f"{{{NAMESPACES['style']}}}name"
	draw_master = f"{{{NAMESPACES['draw']}}}master-page-name"
	page_layout_name = f"{{{NAMESPACES['style']}}}page-layout-name"
	masters = styles_root.findall(".//style:master-page", NAMESPACES)
	master = next(candidate for candidate in masters
		if candidate.attrib[style_name] == page.attrib[draw_master])
	layouts = styles_root.findall(".//style:page-layout", NAMESPACES)
	page_layout = next(candidate for candidate in layouts
		if candidate.attrib[style_name] == master.attrib[page_layout_name])
	properties = page_layout.find("style:page-layout-properties", NAMESPACES)
	require(properties is not None, "ODP page layout defines page dimensions")
	return (parse_odf_length(properties.attrib[f"{{{NAMESPACES['fo']}}}page-width"]),
		parse_odf_length(properties.attrib[f"{{{NAMESPACES['fo']}}}page-height"]))


#============================================
def inspect_odp(odp_path: pathlib.Path) -> None:
	"""Verify editable ODP objects survive the LibreOffice handoff."""
	with zipfile.ZipFile(odp_path) as archive:
		content_root = defusedxml.ElementTree.fromstring(archive.read("content.xml"))
		styles_root = defusedxml.ElementTree.fromstring(archive.read("styles.xml"))
	pages = content_root.findall("./office:body/office:presentation/draw:page", NAMESPACES)
	require(len(pages) == len(LAYOUTS),
		f"ODP has {len(pages)} pages; expected {len(LAYOUTS)} layouts")
	all_text = "\n".join(text_content(page) for page in pages)
	require("title-content acceptance" in all_text and "gallery acceptance" in all_text,
		"ODP preserves editable text across the layout catalog")
	require(any(page.findall(".//text:list-item", NAMESPACES) for page in pages),
		"ODP preserves editable list objects")
	links = [link for page in pages for link in page.findall(".//text:a", NAMESPACES)]
	require(any(link.get(f"{{{NAMESPACES['xlink']}}}href") ==
		"https://example.edu/title-content" for link in links),
		"ODP preserves authored hyperlink targets")
	notes_text = text_content(pages[3].find("presentation:notes", NAMESPACES))
	require("Present this native title-content slide." in notes_text,
		"ODP preserves presenter-note text")
	image_frames = []
	text_boxes = []
	for page in pages:
		page_width, page_height = page_dimensions(styles_root, page)
		for frame in page.findall(".//draw:frame", NAMESPACES):
			if frame.find(".//draw:image", NAMESPACES) is not None:
				image_frames.append(frame)
				width = parse_odf_length(frame.attrib[f"{{{NAMESPACES['svg']}}}width"])
				height = parse_odf_length(frame.attrib[f"{{{NAMESPACES['svg']}}}height"])
				require(width < page_width * MAX_COMPONENT_PAGE_FRACTION or
					height < page_height * MAX_COMPONENT_PAGE_FRACTION,
					"ODP contains no full-slide image frame fallback")
				require(any(text_content(description).startswith("Native ") for description in
					frame.findall(".//svg:desc", NAMESPACES)),
					"ODP component image frame preserves authored description")
			if frame.find(".//draw:text-box", NAMESPACES) is not None:
				text_boxes.append(frame)
	require(len(image_frames) >= 6, "ODP retains component images as separate draw frames")
	require(text_boxes, "ODP retains editable text as separate draw text boxes")


#============================================
def inspect_pdf(pdf_path: pathlib.Path) -> None:
	"""Inspect only the final ODP-derived PDF for its expected page count."""
	result = subprocess.run(["pdfinfo", str(pdf_path)], check=True, capture_output=True, text=True)
	page_match = re.search(r"^Pages:\s+(\d+)\s*$", result.stdout, re.MULTILINE)
	require(page_match is not None, "ODP-derived PDF reports its page count through pdfinfo")
	require(int(page_match.group(1)) == len(LAYOUTS),
		f"ODP-derived PDF has {page_match.group(1)} pages; expected {len(LAYOUTS)}")


#============================================
def run() -> None:
	"""Run the public non-browser native export path for the full layout catalog."""
	root = repo_root()
	output_root = root / "output"
	output_root.mkdir(exist_ok=True)
	stem = f"all_native_layouts_{uuid.uuid4().hex}"
	workspace = pathlib.Path(tempfile.mkdtemp(prefix=f"{stem}_", dir=output_root))
	deck_path = workspace / f"{stem}.md"
	pptx_path = output_root / "pptx" / f"{stem}.pptx"
	odp_path = output_root / "odp" / f"{stem}.odp"
	pdf_path = output_root / "pdf" / f"{stem}.pdf"
	try:
		write_component_images(workspace)
		write_deck(deck_path)
		command = [sys.executable, "tools/marp_export.py", str(deck_path), "--format", "pdf"]
		subprocess.run(command, cwd=root, check=True)
		require(pptx_path.is_file() and odp_path.is_file() and pdf_path.is_file(),
			"Public PDF export creates PPTX, then ODP, then ODP-derived PDF artifacts")
		inspect_pptx(pptx_path)
		inspect_odp(odp_path)
		inspect_pdf(pdf_path)
		print("PASS: all native layouts retain editable PPTX and ODP semantics through PDF")
	finally:
		for artifact_path in (pptx_path, odp_path, pdf_path):
			artifact_path.unlink(missing_ok=True)
		shutil.rmtree(workspace, ignore_errors=True)


if __name__ == "__main__":
	run()
