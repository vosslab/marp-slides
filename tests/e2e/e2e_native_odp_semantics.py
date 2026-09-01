#!/usr/bin/env python3
"""Verify native Marp export semantics through the editable ODP handoff."""

# Standard Library
import pathlib
import re
import shutil
import subprocess
import sys
import tempfile
import uuid
import zipfile

# PIP3 modules
import defusedxml.ElementTree
import PIL.Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


NAMESPACES = {
	"draw": "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0",
	"fo": "urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0",
	"office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
	"presentation": "urn:oasis:names:tc:opendocument:xmlns:presentation:1.0",
	"svg": "urn:oasis:names:tc:opendocument:xmlns:svg-compatible:1.0",
	"style": "urn:oasis:names:tc:opendocument:xmlns:style:1.0",
	"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
	"xlink": "http://www.w3.org/1999/xlink",
}

# Two percent allows writer rounding while rejecting a full-slide raster frame.
MAX_COMPONENT_PAGE_FRACTION = 0.98
ODF_LENGTH_TO_CENTIMETERS = {
	"cm": 1.0,
	"mm": 0.1,
	"in": 2.54,
	"pt": 2.54 / 72.0,
	"pc": 2.54 / 6.0,
	"px": 2.54 / 96.0,
}


#============================================
def repo_root() -> pathlib.Path:
	"""Return the repository root containing this E2E runner."""
	return pathlib.Path(__file__).resolve().parents[2]


#============================================
def require(condition: bool, message: str) -> None:
	"""Raise a diagnostic error when one export contract is not present."""
	if not condition:
		raise RuntimeError(message)


#============================================
def write_component_image(image_path: pathlib.Path) -> None:
	"""Create a small local component image for the Marp deck."""
	image = PIL.Image.new("RGB", (120, 60), (36, 87, 143))
	image.save(image_path)


#============================================
def write_deck(deck_path: pathlib.Path) -> None:
	"""Create a three-slide canonical deck that covers editable components."""
	deck_path.write_text(
		"---\nmarp: true\ntheme: genetics\nsize: 16:10\n---\n"
		"<!-- _class: title-content -->\n"
		"# Native semantics\n\n"
		"- Linked [resource](https://example.edu/native-export)\n"
		"  - Nested editable detail\n"
		"1. Ordered editable step\n\n"
		"<!-- notes: Explain the editable component image. -->\n"
		"---\n"
		"<!-- _class: title-content -->\n"
		"# Component image\n\n"
		"![Native component image: blue calibration swatch](component.png)\n"
		"---\n"
		"<!-- _class: centered-text -->\n"
		"# Second native slide\n",
		encoding="utf-8",
	)


#============================================
def inspect_pptx(pptx_path: pathlib.Path) -> None:
	"""Confirm native PPTX objects preserve every authored component."""
	presentation = Presentation(pptx_path)
	require(len(presentation.slides) == 3, "PPTX contains the expected three native slides")
	first_slide = presentation.slides[0]
	shape_xml = "".join(shape.element.xml for shape in first_slide.shapes)
	shape_text = "\n".join(
		shape.text for shape in first_slide.shapes if shape.has_text_frame
	)
	component_slide = presentation.slides[1]
	pictures = [
		shape for shape in component_slide.shapes
		if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
	]
	require("Native semantics" in shape_text, "PPTX includes editable title text")
	require("Nested editable detail" in shape_text, "PPTX includes editable nested list text")
	require("buChar" in shape_xml and "buAutoNum" in shape_xml,
		"PPTX includes native bullet and ordered-list structure")
	require("hlinkClick" in shape_xml, "PPTX includes a native hyperlink relationship")
	require(first_slide.notes_slide.notes_text_frame.text ==
		"Explain the editable component image.", "PPTX includes presenter-note text")
	require(len(pictures) == 1, "PPTX keeps the component image as one picture object")
	require(pictures[0].width < presentation.slide_width and
		pictures[0].height < presentation.slide_height,
		"PPTX component picture fits inside the slide canvas")


#============================================
def text_content(element: object) -> str:
	"""Return all text nodes in one ODF element."""
	return "".join(element.itertext())


#============================================
def parse_odf_length(length_text: str) -> float:
	"""Return an ODF length in centimeters for supported physical units.

	Args:
		length_text: ODF length text using cm, mm, in, pt, pc, or px.

	Returns:
		Length in centimeters.
	"""
	match = re.fullmatch(r"([+-]?(?:\d+(?:\.\d*)?|\.\d+))(cm|mm|in|pt|pc|px)", length_text)
	require(match is not None,
		f"ODP length {length_text!r} uses one of cm, mm, in, pt, pc, or px")
	value = float(match.group(1))
	require(value >= 0.0, f"ODP length {length_text!r} is non-negative")
	unit = match.group(2)
	centimeters = value * ODF_LENGTH_TO_CENTIMETERS[unit]
	return centimeters


#============================================
def page_dimensions(styles_root: object, page: object) -> tuple[float, float]:
	"""Return the named ODP master page dimensions in centimeters.

	Args:
		styles_root: Parsed ODP styles.xml root element.
		page: Parsed draw:page element from content.xml.

	Returns:
		Page width and height in centimeters.
	"""
	draw_master_page_name = f"{{{NAMESPACES['draw']}}}master-page-name"
	style_name = f"{{{NAMESPACES['style']}}}name"
	style_page_layout_name = f"{{{NAMESPACES['style']}}}page-layout-name"
	fo_page_width = f"{{{NAMESPACES['fo']}}}page-width"
	fo_page_height = f"{{{NAMESPACES['fo']}}}page-height"
	master_page_name = page.attrib[draw_master_page_name]
	master_pages = styles_root.findall(".//style:master-page", NAMESPACES)
	master_page = next(
		candidate for candidate in master_pages
		if candidate.attrib[style_name] == master_page_name
	)
	page_layout_name = master_page.attrib[style_page_layout_name]
	page_layouts = styles_root.findall(".//style:page-layout", NAMESPACES)
	page_layout = next(
		candidate for candidate in page_layouts
		if candidate.attrib[style_name] == page_layout_name
	)
	properties = page_layout.find("style:page-layout-properties", NAMESPACES)
	require(properties is not None, "ODP page layout defines page dimensions")
	page_width = parse_odf_length(properties.attrib[fo_page_width])
	page_height = parse_odf_length(properties.attrib[fo_page_height])
	return page_width, page_height


#============================================
def inspect_odp(odp_path: pathlib.Path) -> None:
	"""Confirm ODP retains editable text, lists, links, notes, and image frames."""
	with zipfile.ZipFile(odp_path) as archive:
		content_root = defusedxml.ElementTree.fromstring(archive.read("content.xml"))
		styles_root = defusedxml.ElementTree.fromstring(archive.read("styles.xml"))
	pages = content_root.findall("./office:body/office:presentation/draw:page", NAMESPACES)
	require(len(pages) == 3, "ODP contains the expected three editable slides")
	first_page = pages[0]
	page_text = text_content(first_page)
	require("Native semantics" in page_text, "ODP includes editable title text")
	require("Nested editable detail" in page_text, "ODP includes editable nested list text")
	require(first_page.findall(".//text:list", NAMESPACES),
		"ODP includes native list containers")
	require(first_page.findall(".//text:list-item", NAMESPACES),
		"ODP includes native list items")
	links = first_page.findall(".//text:a", NAMESPACES)
	require(any(link.get(f"{{{NAMESPACES['xlink']}}}href") ==
		"https://example.edu/native-export" for link in links),
		"ODP includes the authored hyperlink target")
	notes_text = text_content(first_page.find("presentation:notes", NAMESPACES))
	require("Explain the editable component image." in notes_text,
		"ODP includes presenter-note text")
	component_page = pages[1]
	page_width, page_height = page_dimensions(styles_root, component_page)
	frames = component_page.findall(".//draw:frame", NAMESPACES)
	image_frames = [
		frame for frame in frames if frame.find(".//draw:image", NAMESPACES) is not None
	]
	require(len(image_frames) == 1, "ODP keeps the component image in one frame")
	for frame in image_frames:
		width = parse_odf_length(frame.attrib[f"{{{NAMESPACES['svg']}}}width"])
		height = parse_odf_length(frame.attrib[f"{{{NAMESPACES['svg']}}}height"])
		is_full_slide = (width >= page_width * MAX_COMPONENT_PAGE_FRACTION and
			height >= page_height * MAX_COMPONENT_PAGE_FRACTION)
		require(not is_full_slide,
			"ODP component image frame occupies essentially the complete slide canvas")
		descriptions = frame.findall(".//svg:desc", NAMESPACES)
		require(any(text_content(description) ==
			"Native component image: blue calibration swatch" for description in descriptions),
			"ODP preserves the authored component-image description")


#============================================
def run() -> None:
	"""Export a temporary canonical deck and inspect its two editable formats."""
	root = repo_root()
	output_root = root / "output"
	output_root.mkdir(exist_ok=True)
	stem = f"native_odp_e2e_{uuid.uuid4().hex}"
	workspace = pathlib.Path(tempfile.mkdtemp(prefix=f"{stem}_", dir=output_root))
	deck_path = workspace / f"{stem}.md"
	pptx_path = output_root / "pptx" / f"{stem}.pptx"
	odp_path = output_root / "odp" / f"{stem}.odp"
	try:
		write_component_image(workspace / "component.png")
		write_deck(deck_path)
		require(not pptx_path.exists() and not odp_path.exists(),
			"E2E artifact names are available before export")
		command = [sys.executable, "tools/marp_export.py", str(deck_path), "--format", "all"]
		subprocess.run(command, cwd=root, check=True)
		require(pptx_path.is_file() and odp_path.is_file(),
			"Native export creates PPTX and ODP artifacts")
		inspect_pptx(pptx_path)
		inspect_odp(odp_path)
		print("PASS: native PPTX-to-ODP semantic preservation")
	finally:
		for artifact_path in (pptx_path, odp_path, output_root / "pdf" / f"{stem}.pdf"):
			artifact_path.unlink(missing_ok=True)
		shutil.rmtree(workspace, ignore_errors=True)


if __name__ == "__main__":
	run()
