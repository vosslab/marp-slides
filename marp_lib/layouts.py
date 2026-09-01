"""Native editable PPTX builders for the repository's slide-layout vocabulary."""

# Standard Library
import pathlib
from collections.abc import Callable
from dataclasses import dataclass

# PIP3 modules
import PIL.Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

# Local Modules
import marp_lib.native_model


PX = 9525
SLIDE_WIDTH = 1280
SLIDE_HEIGHT = 800
LEFT = 60.0
RIGHT = 1220.0
TITLE_TOP = 52.0
CONTENT_BOTTOM = 754.0
CELL_GUTTER = 42.0
GRID_GUTTER = 24.0
CSS_TO_OFFICE_POINTS = 0.75
BODY_LINE_HEIGHT = 1.3
LIST_ITEM_SPACE_EM = 0.25
MIN_READABLE_BODY_SIZE = 14.0
FONT_NAME = "OpenDyslexic"
URL_FONT_NAME = "PT Sans Narrow"
ACCENT = RGBColor(0x24, 0x57, 0x8F)
FOREGROUND = RGBColor(0x17, 0x20, 0x33)
MUTED = RGBColor(0x52, 0x61, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


class LayoutError(ValueError):
	"""Report an expected source or native-layout validation failure."""


@dataclass(frozen=True)
class LayoutSpec:
	"""Stable authoring contract and renderer for one named slide layout."""
	name: str
	cell_count: int
	allows_root_body: bool
	vertical_title: bool
	vertical_cells: frozenset[int]
	builder: Callable[[object, object, object, "LayoutSpec"], None]


#============================================
def px(value: float) -> Emu:
	"""Convert a 1280x800 CSS-pixel coordinate to Office EMUs."""
	return Emu(round(value * PX))


#============================================
def css_px_to_pt(value: float) -> float:
	"""Convert CSS font pixels to Office points exactly once."""
	return value * CSS_TO_OFFICE_POINTS


#============================================
def inline_text(inlines: tuple[marp_lib.native_model.Inline, ...]) -> str:
	"""Return visible authored text while retaining native runs for rendering."""
	parts: list[str] = []
	for inline in inlines:
		if isinstance(inline, (marp_lib.native_model.Text, marp_lib.native_model.InlineCode)):
			parts.append(inline.value)
		elif isinstance(inline, marp_lib.native_model.Break):
			parts.append(" ")
		else:
			parts.append(inline_text(inline.children))
	return "".join(parts).strip()


#============================================
def add_textbox(slide: object, left: float, top: float, width: float, height: float,
		vertical_anchor: object = MSO_ANCHOR.TOP, vertical_text: bool = False) -> object:
	"""Add one editable native text frame with Office autofit protection."""
	shape = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
	frame = shape.text_frame
	frame.clear()
	frame.margin_left = frame.margin_right = 0
	frame.margin_top = frame.margin_bottom = 0
	frame.vertical_anchor = vertical_anchor
	frame.word_wrap = True
	body_properties = frame._txBody.bodyPr
	for child in list(body_properties):
		if child.tag.endswith(("noAutofit", "normAutofit", "spAutoFit")):
			body_properties.remove(child)
	autofit = OxmlElement("a:normAutofit")
	autofit.set("fontScale", "100000")
	autofit.set("lnSpcReduction", "0")
	body_properties.append(autofit)
	if vertical_text:
		body_properties.set("vert", "vert")
	return frame


#============================================
def write_run(run: object, text: str, size: float, color: object, bold: bool = False,
		italic: bool = False, url: str | None = None, displayed_url: bool = False) -> None:
	"""Apply the repository font contract to one editable run."""
	run.text = text
	run.font.name = URL_FONT_NAME if displayed_url else FONT_NAME
	run.font.size = Pt(css_px_to_pt(size))
	run.font.bold = bold
	run.font.italic = italic
	run.font.color.rgb = color
	if url is not None:
		run.hyperlink.address = url
		run.font.underline = True


#============================================
def add_inline_runs(paragraph: object, inlines: tuple[marp_lib.native_model.Inline, ...], size: float,
		color: object = FOREGROUND, bold: bool = False, italic: bool = False,
		url: str | None = None, displayed_url: bool = False) -> None:
	"""Write the typed inline tree as editable, formatted Office text runs."""
	for inline in inlines:
		if isinstance(inline, marp_lib.native_model.Text):
			write_run(paragraph.add_run(), inline.value, size, color, bold, italic, url, displayed_url)
		elif isinstance(inline, marp_lib.native_model.InlineCode):
			write_run(paragraph.add_run(), inline.value, size, color, bold, italic, url, displayed_url)
		elif isinstance(inline, marp_lib.native_model.Break):
			paragraph.add_line_break()
		elif isinstance(inline, marp_lib.native_model.Strong):
			add_inline_runs(paragraph, inline.children, size, color, True, italic, url, displayed_url)
		elif isinstance(inline, marp_lib.native_model.Emphasis):
			add_inline_runs(paragraph, inline.children, size, color, bold, True, url, displayed_url)
		elif isinstance(inline, marp_lib.native_model.Link):
			is_literal_url = inline_text(inline.children) == inline.url
			add_inline_runs(paragraph, inline.children, size, ACCENT, bold, italic, inline.url, is_literal_url)
	if not paragraph.runs:
		write_run(paragraph.add_run(), "", size, color, bold, italic, url, displayed_url)


#============================================
def flatten_list(block: marp_lib.native_model.ListBlock, level: int = 0) -> list[tuple[tuple[marp_lib.native_model.Inline, ...], int, bool, bool, int]]:
	"""Flatten nested lists into Office paragraphs, preserving ordered starts."""
	items: list[tuple[tuple[marp_lib.native_model.Inline, ...], int, bool, bool, int]] = []
	for offset, item in enumerate(block.items):
		items.append((item.inlines, level, block.ordered, False, block.start + offset))
		for child in item.children:
			items.extend(flatten_list(child, level + 1))
	return items


#============================================
def body_parts(blocks: tuple[marp_lib.native_model.Block, ...]) -> tuple[list[marp_lib.native_model.Heading], list[tuple[tuple[marp_lib.native_model.Inline, ...], int, bool, bool, int]], list[marp_lib.native_model.Image]]:
	"""Classify already typed blocks without reparsing canonical Markdown."""
	headings: list[marp_lib.native_model.Heading] = []
	items: list[tuple[tuple[marp_lib.native_model.Inline, ...], int, bool, bool, int]] = []
	images: list[marp_lib.native_model.Image] = []
	for block in blocks:
		if isinstance(block, marp_lib.native_model.Heading):
			headings.append(block)
		elif isinstance(block, marp_lib.native_model.Paragraph):
			items.append((block.inlines, 0, False, True, 1))
		elif isinstance(block, marp_lib.native_model.ListBlock):
			items.extend(flatten_list(block))
		else:
			images.append(block)
	return headings, items, images


#============================================
def wrapped_line_count(inlines: tuple[marp_lib.native_model.Inline, ...], size: float, width: float,
		level: int = 0) -> int:
	"""Estimate wrapped editable text lines conservatively."""
	available_width = max(width - level * size * 1.43, size * 3)
	characters_per_line = max(int(available_width / (size * 0.54)), 1)
	words = inline_text(inlines).split()
	if not words:
		return 1
	lines, current = 1, 0
	for word in words:
		if current and current + len(word) + 1 > characters_per_line:
			lines, current = lines + 1, len(word)
		else:
			current += len(word) + (1 if current else 0)
		lines += max((len(word) - 1) // characters_per_line, 0)
	return lines


#============================================
def estimate_items_height(items: list[tuple[tuple[marp_lib.native_model.Inline, ...], int, bool, bool, int]], size: float, width: float) -> float:
	"""Estimate native paragraph height in CSS pixels."""
	return sum(wrapped_line_count(inlines, size, width, level) * size * BODY_LINE_HEIGHT +
		size * LIST_ITEM_SPACE_EM for inlines, level, _, _, _ in items)


#============================================
def fit_body_size(item_sets: list[list[tuple[tuple[marp_lib.native_model.Inline, ...], int, bool, bool, int]]], width: float,
		height: float, preferred_size: float, source: marp_lib.native_model.SourceLocation,
		context: str) -> float:
	"""Choose one shared readable body size or report a capacity violation."""
	for quarter_points in range(int(preferred_size * 4), int(MIN_READABLE_BODY_SIZE * 4) - 1, -1):
		size = quarter_points / 4
		if all(estimate_items_height(items, size, width) <= height for items in item_sets):
			return size
	raise LayoutError(f"{source.path}:{source.line}: {context} content cannot fit within the supported readable minimum of "
		f"{MIN_READABLE_BODY_SIZE:g} CSS px")


#============================================
def write_items(frame: object, items: list[tuple[tuple[marp_lib.native_model.Inline, ...], int, bool, bool, int]], size: float,
		first_paragraph: object | None = None) -> None:
	"""Write paragraphs using native bullet and automatic-number OOXML."""
	for index, (inlines, level, ordered, paragraph_only, start) in enumerate(items):
		paragraph = first_paragraph if index == 0 and first_paragraph is not None else (
			frame.paragraphs[0] if index == 0 else frame.add_paragraph())
		paragraph.level = level
		paragraph.space_after = Pt(css_px_to_pt(size * LIST_ITEM_SPACE_EM))
		paragraph.line_spacing = BODY_LINE_HEIGHT
		bullet = OxmlElement("a:buNone" if paragraph_only else
			"a:buAutoNum" if ordered else "a:buChar")
		if ordered:
			bullet.set("type", "arabicPeriod")
			bullet.set("startAt", str(start))
		elif not paragraph_only:
			bullet.set("char", "\u2022")
		paragraph._p.get_or_add_pPr().insert(0, bullet)
		add_inline_runs(paragraph, inlines, size)


#============================================
def add_background(slide: object) -> None:
	"""Write the white canvas and native blue accent rule."""
	slide.background.fill.solid()
	slide.background.fill.fore_color.rgb = WHITE
	accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(0), px(0), px(SLIDE_WIDTH), px(10))
	accent.fill.solid()
	accent.fill.fore_color.rgb = ACCENT
	accent.line.fill.background()


#============================================
def resolve_image_path(deck: object, image: marp_lib.native_model.Image) -> pathlib.Path:
	"""Resolve a component image within its canonical repository boundary."""
	image_path = (deck.asset_root / image.source).resolve()
	if not image_path.is_relative_to(deck.repo_root):
		raise LayoutError(f"{image.location.path}:{image.location.line}: component image must be inside the repository: {image.source}")
	if not image_path.is_file():
		raise LayoutError(f"{image.location.path}:{image.location.line}: component image is missing: {image.source}")
	return image_path


#============================================
def add_picture(slide: object, image_path: pathlib.Path, image: marp_lib.native_model.Image,
		left: float, top: float, width: float, height: float) -> None:
	"""Add one contained component picture with its authored description."""
	with PIL.Image.open(image_path) as opened_image:
		image_width, image_height = opened_image.size
	scale = min(width / image_width, height / image_height)
	display_width, display_height = image_width * scale, image_height * scale
	if display_width >= SLIDE_WIDTH or display_height >= SLIDE_HEIGHT:
		raise LayoutError(f"{image.location.path}:{image.location.line}: component image occupies a full slide: {image.source}")
	picture = slide.shapes.add_picture(str(image_path), px(left + (width - display_width) / 2),
		px(top + (height - display_height) / 2), px(display_width), px(display_height))
	picture.element.nvPicPr.cNvPr.set("descr", image.alt_text)


#============================================
def title_size(source: marp_lib.native_model.Slide, default_size: float) -> float:
	"""Return the typed H1 override or the layout's established default size."""
	return float(source.title_size_override.preset.value) if source.title_size_override else default_size


#============================================
def require_title_capacity(source: marp_lib.native_model.Slide, title: marp_lib.native_model.Heading,
		size: float, width: float, height: float, layout: str, vertical: bool = False) -> float:
	"""Measure an explicit H1 request without silently shrinking authored type."""
	if vertical:
		line_count = max(len(inline_text(title.inlines).replace(" ", "")), 1)
	else:
		line_count = wrapped_line_count(title.inlines, size, width)
	required_height = line_count * size * 1.12
	if source.title_size_override and required_height > height:
		preset = source.title_size_override.preset.value
		raise layout_error(title.location,
			f"{layout} H1 font-size-{preset} does not fit its native title region")
	return required_height


#============================================
def title_and_content_top(slide: object, source: marp_lib.native_model.Slide,
		title: marp_lib.native_model.Heading, vertical_title: bool = False) -> tuple[float, float, float]:
	"""Write a title and return the remaining content rectangle origin and size."""
	if vertical_title:
		frame = add_textbox(slide, LEFT, 60, 94, 666, vertical_text=True)
		paragraph = frame.paragraphs[0]
		size = title_size(source, 38)
		require_title_capacity(source, title, size, 94, 666, source.layout_class, True)
		add_inline_runs(paragraph, title.inlines, size)
		for run in paragraph.runs:
			run.font.bold = True
		return 178, 82, RIGHT - 178
	size = title_size(source, 48)
	title_height = require_title_capacity(source, title, size, RIGHT - LEFT, 170, source.layout_class)
	frame = add_textbox(slide, LEFT, TITLE_TOP, RIGHT - LEFT, title_height)
	paragraph = frame.paragraphs[0]
	add_inline_runs(paragraph, title.inlines, size)
	for run in paragraph.runs:
		run.font.bold = True
	return LEFT, TITLE_TOP + title_height + 24, RIGHT - LEFT


#============================================
def layout_error(source: marp_lib.native_model.SourceLocation | marp_lib.native_model.Slide,
		message: str) -> ValueError:
	"""Attach a layout-capacity failure to its canonical source location."""
	location = source.location if isinstance(source, marp_lib.native_model.Slide) else source
	return ValueError(f"{location.path}:{location.line}: {message}")


#============================================
def validate_layout_source(source: marp_lib.native_model.Slide) -> LayoutSpec:
	"""Select one layout and prove every supported source block has a destination."""
	spec = LAYOUTS[source.layout_class]
	headings, items, images = body_parts(source.blocks)
	cells = source.cells
	if spec.name == "blank":
		if source.blocks or cells:
			offending = source.blocks[0].location if source.blocks else cells[0].location
			raise layout_error(offending, "blank slides must be empty")
		return spec
	if spec.cell_count:
		if len(cells) != spec.cell_count:
			offending = cells[-1].location if cells else source.location
			raise layout_error(offending, f"{spec.name} slides require exactly {spec.cell_count} blockquote cells")
		if items or images or any(heading.level != 1 for heading in headings):
			offending = next((block.location for block in source.blocks if not isinstance(block,
				marp_lib.native_model.Heading) or block.level != 1), source.location)
			raise layout_error(offending, f"{spec.name} slides support a title and blockquote cells only")
		if len(headings) != 1 or headings[0].level != 1:
			offending = headings[1].location if len(headings) > 1 else source.location
			raise layout_error(offending, f"{spec.name} slides require exactly one level-one title")
		for index, cell in enumerate(cells, start=1):
			cell_headings, cell_items, cell_images = body_parts(cell.blocks)
			if len(cell_headings) > 1 or any(heading.level != 2 for heading in cell_headings):
				offending = next((block.location for block in cell.blocks if isinstance(block,
					marp_lib.native_model.Heading) and block.level != 2), cell.location)
				raise layout_error(offending, f"{spec.name} cell {index} supports one optional level-two heading")
			if cell_items and cell_images:
				offending = next(block.location for block in cell.blocks if isinstance(block,
					marp_lib.native_model.Image))
				raise layout_error(offending, f"{spec.name} cell {index} cannot combine text and component images")
			if not cell_items and not cell_images:
				raise layout_error(cell.location, f"{spec.name} cell {index} requires editable text or component images")
		return spec
	if cells:
		raise layout_error(cells[0].location, f"{spec.name} slides do not accept blockquote cells")
	if spec.name in ("title-slide", "centered-text"):
		if not headings or headings[0].level != 1 or items or images:
			raise layout_error(source, f"{spec.name} slides support a title and level-two subtitle lines only")
		if any(heading.level != 2 for heading in headings[1:]):
			raise layout_error(source, f"{spec.name} subtitle lines must use level-two Markdown")
		return spec
	if spec.name == "title-only":
		if len(headings) != 1 or headings[0].level != 1 or items or images:
			raise layout_error(source, "title-only slides require exactly one level-one title")
		return spec
	if spec.name == "gallery":
		if len(headings) > 1 or any(heading.level != 1 for heading in headings) or items or not 2 <= len(images) <= 6:
			raise layout_error(source, "gallery slides support an optional title and two through six component images")
		return spec
	if len(headings) != 1 or headings[0].level != 1:
		offending = headings[1].location if len(headings) > 1 else source.location
		raise layout_error(offending, f"{spec.name} slides require exactly one level-one title")
	if spec.name in ("title-vertical-text", "vertical-title-vertical-text"):
		body_blocks = source.blocks[1:]
		if len(body_blocks) != 1:
			offending = body_blocks[1].location if len(body_blocks) > 1 else headings[0].location
			raise layout_error(offending, f"{spec.name} slides require exactly one root body block")
		if not isinstance(body_blocks[0], (marp_lib.native_model.Paragraph,
				marp_lib.native_model.ListBlock, marp_lib.native_model.Image)):
			raise layout_error(body_blocks[0].location,
				f"{spec.name} slides require exactly one root body block")
	if not spec.allows_root_body or (items and images):
		offending = next((block.location for block in source.blocks[1:] if isinstance(block,
			marp_lib.native_model.Image)), headings[0].location)
		raise layout_error(offending, f"{spec.name} slides require one body mode: editable text or component images")
	if not items and not images:
		raise layout_error(headings[0].location, f"{spec.name} slides require editable body text or one component image")
	if len(images) > 1:
		raise layout_error(images[1].location, f"{spec.name} slides support one contained component image")
	return spec


def render_cell(slide: object, deck: marp_lib.native_model.Deck, cell: marp_lib.native_model.Cell, rectangle: tuple[float, float, float, float],
		vertical: bool = False) -> None:
	"""Render one independently editable cell in its assigned rectangle."""
	left, top, width, height = rectangle
	headings, items, images = body_parts(cell.blocks)
	body_top, body_height = top, height
	if headings:
		heading = headings[0]
		heading_size = 28
		heading_height = wrapped_line_count(heading.inlines, heading_size, width) * heading_size * 1.2
		head_frame = add_textbox(slide, left, top, width, heading_height, vertical_text=vertical)
		heading_paragraph = head_frame.paragraphs[0]
		add_inline_runs(heading_paragraph, heading.inlines, heading_size)
		for run in heading_paragraph.runs:
			run.font.bold = True
		body_top, body_height = top + heading_height + 10, height - heading_height - 10
	if images:
		gap = 12
		image_width = (width - gap * (len(images) - 1)) / len(images)
		for index, image in enumerate(images):
			add_picture(slide, resolve_image_path(deck, image), image,
				left + index * (image_width + gap), body_top, image_width, body_height)
	elif items:
		body_block = next(block for block in cell.blocks if not isinstance(block,
			marp_lib.native_model.Heading))
		size = fit_body_size([items], width, body_height, 22, body_block.location, "cell")
		frame = add_textbox(slide, left, body_top, width, body_height, vertical_text=vertical)
		write_items(frame, items, size)


#============================================
def grid_rectangles(left: float, top: float, width: float, height: float, columns: int,
		rows: int) -> list[tuple[float, float, float, float]]:
	"""Return reading-order grid cells from shared geometry constants."""
	cell_width = (width - CELL_GUTTER * (columns - 1)) / columns
	cell_height = (height - GRID_GUTTER * (rows - 1)) / rows
	return [(left + column * (cell_width + CELL_GUTTER), top + row * (cell_height + GRID_GUTTER),
		cell_width, cell_height) for row in range(rows) for column in range(columns)]


#============================================
def build_blank(slide: object, source: object, deck: object, spec: LayoutSpec) -> None:
	"""Render an empty native canvas."""


#============================================
def build_title_slide(slide: object, source: marp_lib.native_model.Slide, deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render centered title and optional subtitle."""
	headings, _, _ = body_parts(source.blocks)
	frame = add_textbox(slide, 110, 180, 1060, 390, MSO_ANCHOR.MIDDLE)
	title_size_value = title_size(source, 60)
	subtitle_height = max(len(headings) - 1, 0) * 31 * 1.2
	require_title_capacity(source, headings[0], title_size_value, 1060, 390 - subtitle_height,
		source.layout_class)
	for index, heading in enumerate(headings):
		paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
		paragraph.alignment = PP_ALIGN.CENTER
		paragraph.space_after = Pt(14)
		add_inline_runs(paragraph, heading.inlines, title_size_value if index == 0 else 31,
			FOREGROUND if index == 0 else MUTED)
		for run in paragraph.runs:
			run.font.bold = index == 0


#============================================
def build_title_only(slide: object, source: marp_lib.native_model.Slide, deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render only the native title region."""
	title = body_parts(source.blocks)[0][0]
	if source.title_size_override is None:
		title_and_content_top(slide, source, title)
		return
	size = title_size(source, 48)
	require_title_capacity(source, title, size, 1060, 650, source.layout_class)
	frame = add_textbox(slide, 110, 76, 1060, 650, MSO_ANCHOR.MIDDLE)
	paragraph = frame.paragraphs[0]
	paragraph.alignment = PP_ALIGN.CENTER
	add_inline_runs(paragraph, title.inlines, size)
	for run in paragraph.runs:
		run.font.bold = True


#============================================
def build_centered_text(slide: object, source: marp_lib.native_model.Slide, deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render centered editable title and optional subtitle."""
	build_title_slide(slide, source, deck, spec)


#============================================
def render_root_body(slide: object, source: marp_lib.native_model.Slide, deck: marp_lib.native_model.Deck,
		spec: LayoutSpec) -> None:
	"""Render title and one native editable body or component image region."""
	headings, items, images = body_parts(source.blocks)
	left, top, width = title_and_content_top(slide, source, headings[0], spec.vertical_title)
	height = CONTENT_BOTTOM - top
	if images:
		for image in images:
			add_picture(slide, resolve_image_path(deck, image), image, left, top, width, height)
	elif items:
		body_block = next(block for block in source.blocks if not isinstance(block,
			marp_lib.native_model.Heading))
		size = fit_body_size([items], width, height, 26, body_block.location, spec.name)
		frame = add_textbox(slide, left, top, width, height, vertical_text=bool(spec.vertical_cells))
		write_items(frame, items, size)


#============================================
def build_title_content(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render the ordinary title-and-content native layout."""
	render_root_body(slide, source, deck, spec)


#============================================
def build_vertical_title_vertical_text(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render one vertical body pane beside its vertical title strip."""
	render_root_body(slide, source, deck, spec)


#============================================
def build_title_vertical_text(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render one vertical body pane below its horizontal title."""
	render_root_body(slide, source, deck, spec)


#============================================
def render_cells(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec,
		rectangles: list[tuple[float, float, float, float]]) -> None:
	"""Render already selected native cell rectangles in source reading order."""
	for index, (cell, rectangle) in enumerate(zip(source.cells, rectangles)):
		render_cell(slide, deck, cell, rectangle, index in spec.vertical_cells)


#============================================
def content_rectangle(slide: object, source: marp_lib.native_model.Slide,
		spec: LayoutSpec) -> tuple[float, float, float, float]:
	"""Return the common native body region after emitting its title."""
	title = body_parts(source.blocks)[0][0]
	left, top, width = title_and_content_top(slide, source, title, spec.vertical_title)
	return left, top, width, CONTENT_BOTTOM - top


#============================================
def build_title_two_content(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render the two-peer content layout."""
	left, top, width, height = content_rectangle(slide, source, spec)
	render_cells(slide, source, deck, spec, grid_rectangles(left, top, width, height, 2, 1))


#============================================
def build_title_content_and_two_content(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render one left cell beside two vertically stacked right cells."""
	left, top, width, height = content_rectangle(slide, source, spec)
	cell_width = (width - CELL_GUTTER) / 2
	half_height = (height - GRID_GUTTER) / 2
	rectangles = [(left, top, cell_width, height), (left + cell_width + CELL_GUTTER, top, cell_width, half_height),
		(left + cell_width + CELL_GUTTER, top + half_height + GRID_GUTTER, cell_width, half_height)]
	render_cells(slide, source, deck, spec, rectangles)


#============================================
def build_title_two_content_and_content(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render two vertically stacked left cells beside one right cell."""
	left, top, width, height = content_rectangle(slide, source, spec)
	cell_width = (width - CELL_GUTTER) / 2
	half_height = (height - GRID_GUTTER) / 2
	rectangles = [(left, top, cell_width, half_height), (left, top + half_height + GRID_GUTTER, cell_width, half_height),
		(left + cell_width + CELL_GUTTER, top, cell_width, height)]
	render_cells(slide, source, deck, spec, rectangles)


#============================================
def build_title_content_over_content(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render two vertically stacked full-width content cells."""
	left, top, width, height = content_rectangle(slide, source, spec)
	render_cells(slide, source, deck, spec, grid_rectangles(left, top, width, height, 1, 2))


#============================================
def build_title_two_content_over_content(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render two upper peers above one full-width lower cell."""
	left, top, width, height = content_rectangle(slide, source, spec)
	half_height = (height - GRID_GUTTER) / 2
	rectangles = grid_rectangles(left, top, width, half_height, 2, 1)
	rectangles.append((left, top + half_height + GRID_GUTTER, width, half_height))
	render_cells(slide, source, deck, spec, rectangles)


#============================================
def build_title_four_content(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render a two-by-two native content grid."""
	left, top, width, height = content_rectangle(slide, source, spec)
	render_cells(slide, source, deck, spec, grid_rectangles(left, top, width, height, 2, 2))


#============================================
def build_title_six_content(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render a three-by-two native content grid."""
	left, top, width, height = content_rectangle(slide, source, spec)
	render_cells(slide, source, deck, spec, grid_rectangles(left, top, width, height, 3, 2))


#============================================
def build_vertical_title_text_chart(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render two peer cells beside the fixed-width vertical title strip."""
	left, top, width, height = content_rectangle(slide, source, spec)
	render_cells(slide, source, deck, spec, grid_rectangles(left, top, width, height, 2, 1))


#============================================
def build_title_two_vertical_text_clipart(slide: object, source: marp_lib.native_model.Slide,
		deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render two stacked cells beside one vertical native pane."""
	left, top, width, height = content_rectangle(slide, source, spec)
	right_width = (width - CELL_GUTTER) * .34
	left_width = width - CELL_GUTTER - right_width
	half_height = (height - GRID_GUTTER) / 2
	rectangles = [(left, top, left_width, half_height), (left, top + half_height + GRID_GUTTER, left_width, half_height),
		(left + left_width + CELL_GUTTER, top, right_width, height)]
	render_cells(slide, source, deck, spec, rectangles)


#============================================
def build_gallery(slide: object, source: marp_lib.native_model.Slide, deck: marp_lib.native_model.Deck, spec: LayoutSpec) -> None:
	"""Render a row of independently contained component images."""
	headings, _, images = body_parts(source.blocks)
	if headings:
		_, top, _ = title_and_content_top(slide, source, headings[0])
	else:
		top = 82
	width = (RIGHT - LEFT - 18 * (len(images) - 1)) / len(images)
	for index, image in enumerate(images):
		add_picture(slide, resolve_image_path(deck, image), image, LEFT + index * (width + 18),
			top, width, CONTENT_BOTTOM - top)


LAYOUTS: dict[str, LayoutSpec] = {
	"blank": LayoutSpec("blank", 0, False, False, frozenset(), build_blank),
	"title-only": LayoutSpec("title-only", 0, False, False, frozenset(), build_title_only),
	"title-slide": LayoutSpec("title-slide", 0, False, False, frozenset(), build_title_slide),
	"title-content": LayoutSpec("title-content", 0, True, False, frozenset(), build_title_content),
	"centered-text": LayoutSpec("centered-text", 0, False, False, frozenset(), build_centered_text),
	"title-two-content": LayoutSpec("title-two-content", 2, False, False, frozenset(), build_title_two_content),
	"title-content-and-two-content": LayoutSpec("title-content-and-two-content", 3, False, False, frozenset(), build_title_content_and_two_content),
	"title-two-content-and-content": LayoutSpec("title-two-content-and-content", 3, False, False, frozenset(), build_title_two_content_and_content),
	"title-content-over-content": LayoutSpec("title-content-over-content", 2, False, False, frozenset(), build_title_content_over_content),
	"title-two-content-over-content": LayoutSpec("title-two-content-over-content", 3, False, False, frozenset(), build_title_two_content_over_content),
	"title-four-content": LayoutSpec("title-four-content", 4, False, False, frozenset(), build_title_four_content),
	"title-six-content": LayoutSpec("title-six-content", 6, False, False, frozenset(), build_title_six_content),
	"vertical-title-vertical-text": LayoutSpec("vertical-title-vertical-text", 0, True, True, frozenset({0}), build_vertical_title_vertical_text),
	"vertical-title-text-chart": LayoutSpec("vertical-title-text-chart", 2, False, True, frozenset(), build_vertical_title_text_chart),
	"title-vertical-text": LayoutSpec("title-vertical-text", 0, True, False, frozenset({0}), build_title_vertical_text),
	"title-two-vertical-text-clipart": LayoutSpec("title-two-vertical-text-clipart", 3, False, False, frozenset({2}), build_title_two_vertical_text_clipart),
	"gallery": LayoutSpec("gallery", 0, False, False, frozenset(), build_gallery),
}


#============================================
def render_layout(slide: object, source: marp_lib.native_model.Slide, deck: marp_lib.native_model.Deck) -> None:
	"""Validate and render exactly one native layout on a blank slide."""
	spec = validate_layout_source(source)
	add_background(slide)
	spec.builder(slide, source, deck, spec)
