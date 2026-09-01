"""Export repository Marp Markdown as native editable presentations."""

# Standard Library
import html
import os
import pathlib
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass

# PIP3 modules
import PIL.Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt


PX = 9525
SLIDE_WIDTH = 1280
SLIDE_HEIGHT = 800
CSS_DPI = 96
OFFICE_DPI = 72
CSS_TO_OFFICE_POINTS = OFFICE_DPI / CSS_DPI
BODY_LINE_HEIGHT = 1.3
LIST_ITEM_SPACE_EM = 0.25
MIN_READABLE_BODY_SIZE = 14.0
FONT_NAME = "OpenDyslexic"
URL_FONT_NAME = "PT Sans Narrow"
ACCENT = RGBColor(0x24, 0x57, 0x8F)
FOREGROUND = RGBColor(0x17, 0x20, 0x33)
MUTED = RGBColor(0x52, 0x61, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLASS_NAMES = {"lead", "figure", "gallery", "two-pane", "dense", "list-columns", "url-list"}
IMAGE_PATTERN = re.compile(r"!\[([^]]*)\]\(([^)]+)\)")
HEADING_PATTERN = re.compile(r"^(#{1,2})\s+(.+)$")
LIST_PATTERN = re.compile(r"^(\s*)([-*+] |\d+\. )(.+)$")
COMMENT_PATTERN = re.compile(r"<!--\s*(.*?)\s*-->", re.DOTALL)
BG_IMAGE_PATTERN = re.compile(r"^bg\s+(left|right):(\d+)%\s+contain$")
SOURCE_FALLBACK_PATTERNS = (
	re.compile(r"\bslide_\d+_source\.(?:png|jpe?g|webp)\b", re.IGNORECASE),
	re.compile(r"_class\s*:\s*source-fallback\b", re.IGNORECASE),
)


@dataclass
class Deck:
	"""Parsed source and its asset boundary."""
	asset_root: pathlib.Path
	repo_root: pathlib.Path
	title: str
	slides: list["Slide"]


@dataclass
class Slide:
	"""One supported slide with presenter metadata."""
	content: str
	classes: set[str]
	paginate: bool
	notes: str


#============================================
def px(value: float) -> Emu:
	"""Convert canvas pixels to Office EMUs."""
	converted = Emu(round(value * PX))
	return converted


#============================================
def css_px_to_pt(value: float) -> float:
	"""Convert one CSS pixel measurement to an Office point measurement."""
	converted = value * CSS_TO_OFFICE_POINTS
	return converted


#============================================
def find_repo_root() -> pathlib.Path:
	"""Return the current Git repository root."""
	result = subprocess.run(["git", "rev-parse", "--show-toplevel"], check=True,
		capture_output=True, text=True)
	repo_root = pathlib.Path(result.stdout.strip()).resolve()
	return repo_root


#============================================
def find_source_root(input_path: pathlib.Path) -> pathlib.Path:
	"""Find an input checkout, supporting standalone test fixtures."""
	for candidate in (input_path.parent, *input_path.parents):
		if (candidate / ".git").exists():
			return candidate.resolve()
	return input_path.parent.resolve()


#============================================
def validate_input(input_value: str, repo_root: pathlib.Path) -> pathlib.Path:
	"""Resolve a canonical Markdown deck and reject source-slide fallbacks."""
	input_path = pathlib.Path(input_value).expanduser().resolve()
	if not input_path.is_file():
		raise ValueError(f"input is not a file: {input_value}")
	if not input_path.is_relative_to(repo_root):
		raise ValueError("input must be inside this repository")
	if input_path.suffix != ".md":
		raise ValueError("input must use the .md extension")
	markdown = input_path.read_text(encoding="utf-8")
	if any(pattern.search(markdown) for pattern in SOURCE_FALLBACK_PATTERNS):
		raise ValueError("full-slide source images are failed conversions and cannot be exported")
	return input_path


#============================================
def parse_front_matter(markdown: str) -> tuple[str, bool, str]:
	"""Validate the canonical document contract and return its body settings."""
	if not markdown.startswith("---\n"):
		raise ValueError("Marp Markdown must begin with YAML front matter")
	front_matter_end = markdown.find("\n---", 4)
	if front_matter_end == -1:
		raise ValueError("Marp Markdown front matter is not closed")
	front_matter = markdown[4:front_matter_end]
	if not re.search(r"^marp:\s*true\s*$", front_matter, re.MULTILINE):
		raise ValueError("Markdown front matter must declare marp: true")
	if not re.search(r"^theme:\s*genetics\s*$", front_matter, re.MULTILINE):
		raise ValueError("Markdown front matter must declare theme: genetics")
	if not re.search(r"^size:\s*16:10\s*$", front_matter, re.MULTILINE):
		raise ValueError("Markdown front matter must declare size: 16:10")
	paginate_match = re.search(r"^paginate:\s*(true|false)\s*$", front_matter, re.MULTILINE)
	title_match = re.search(r"^title:\s*(.+?)\s*$", front_matter, re.MULTILINE)
	body = markdown[front_matter_end + 4:].lstrip("\n")
	default_paginate = paginate_match is None or paginate_match.group(1) == "true"
	title = title_match.group(1).strip(" '\"") if title_match else ""
	return body, default_paginate, title


#============================================
def extract_slide(raw_slide: str, default_paginate: bool) -> Slide:
	"""Extract slide directives and comments as presenter notes."""
	classes: set[str] = set()
	paginate = default_paginate
	notes: list[str] = []
	continuations: list[str] = []
	for match in COMMENT_PATTERN.finditer(raw_slide):
		comment = match.group(1).strip()
		if comment.startswith("_class:"):
			classes.update(comment.partition(":")[2].strip().split())
		elif comment.startswith("_paginate:"):
			paginate = comment.partition(":")[2].strip().lower() != "false"
		elif comment.startswith("ODP hidden slide skipped:"):
			continue
		elif comment.startswith("Continuation:"):
			continuations.append(comment)
		elif comment:
			notes.append(comment.removeprefix("notes:").strip())
	unsupported = classes - CLASS_NAMES
	if unsupported:
		raise ValueError(f"unsupported Marp slide class: {' '.join(sorted(unsupported))}")
	content = COMMENT_PATTERN.sub("", raw_slide).strip()
	if re.search(r"<(?!\!--)", content):
		raise ValueError("raw HTML is outside the supported Marp vocabulary")
	if continuations:
		content += "\n" + "\n".join(continuations)
	slide = Slide(content, classes, paginate, "\n\n".join(notes).strip())
	return slide


#============================================
def parse_deck(input_path: pathlib.Path) -> Deck:
	"""Parse canonical Marp Markdown into the native slide vocabulary."""
	markdown = input_path.read_text(encoding="utf-8")
	body, default_paginate, title = parse_front_matter(markdown)
	raw_slides = [item for item in re.split(r"\n---\s*(?:\n|$)", body) if item.strip()]
	if not raw_slides:
		raise ValueError("Marp Markdown contains no slides")
	slides = [extract_slide(raw_slide, default_paginate) for raw_slide in raw_slides]
	deck = Deck(input_path.parent.resolve(), find_source_root(input_path), title, slides)
	return deck


#============================================
def resolve_image_path(deck: Deck, image_value: str) -> pathlib.Path:
	"""Resolve one component image inside its source repository."""
	image_path = (deck.asset_root / image_value).resolve()
	if not image_path.is_relative_to(deck.repo_root):
		raise ValueError(f"component image must be inside the repository: {image_value}")
	if not image_path.is_file():
		raise ValueError(f"component image is missing: {image_value}")
	return image_path


#============================================
def image_bounds(image_path: pathlib.Path, left: float, top: float, width: float,
		height: float) -> tuple[Emu, Emu, Emu, Emu]:
	"""Return contain-sized coordinates inside a layout rectangle."""
	with PIL.Image.open(image_path) as image:
		image_width, image_height = image.size
	scale = min(width / image_width, height / image_height)
	display_width = image_width * scale
	display_height = image_height * scale
	return (
		px(left + (width - display_width) / 2),
		px(top + (height - display_height) / 2),
		px(display_width),
		px(display_height),
	)


#============================================
def add_picture(slide: object, image_path: pathlib.Path, alt_text: str, left: float, top: float,
		width: float, height: float) -> None:
	"""Add a component picture with contain geometry."""
	image_left, image_top, image_width, image_height = image_bounds(
		image_path, left, top, width, height,
	)
	if image_width >= px(SLIDE_WIDTH) or image_height >= px(SLIDE_HEIGHT):
		raise ValueError(f"component image occupies a full slide: {image_path}")
	picture = slide.shapes.add_picture(str(image_path), image_left, image_top, image_width, image_height)
	picture.element.nvPicPr.cNvPr.set("descr", alt_text)


#============================================
def add_textbox(slide: object, left: float, top: float, width: float, height: float,
		vertical: object = MSO_ANCHOR.TOP) -> object:
	"""Add one editable text box using calibrated geometry."""
	shape = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
	text_frame = shape.text_frame
	text_frame.clear()
	text_frame.margin_left = 0
	text_frame.margin_right = 0
	text_frame.margin_top = 0
	text_frame.margin_bottom = 0
	text_frame.vertical_anchor = vertical
	text_frame.word_wrap = True
	# Keep an editable native text-fit instruction as a final safeguard. The
	# layout pass computes the actual shared font size before this is emitted.
	body_properties = text_frame._txBody.bodyPr
	for child in list(body_properties):
		if child.tag.endswith(("noAutofit", "normAutofit", "spAutoFit")):
			body_properties.remove(child)
	autofit = OxmlElement("a:normAutofit")
	autofit.set("fontScale", "100000")
	autofit.set("lnSpcReduction", "0")
	body_properties.append(autofit)
	return text_frame


#============================================
def clean_inline(text: str) -> str:
	"""Convert supported inline Markdown to editable text."""
	cleaned = re.sub(r"!\[[^]]*\]\([^)]+\)", "", text)
	cleaned = re.sub(r"\[([^]]+)\]\([^)]+\)", r"\1", cleaned)
	cleaned = re.sub(r"\*\*(.+?)\*\*", r"\1", cleaned)
	cleaned = re.sub(r"__(.+?)__", r"\1", cleaned)
	cleaned = html.unescape(cleaned)
	return cleaned.strip()


#============================================
def write_run(run: object, text: str, size: float, color: object, bold: bool = False,
		url_font: bool = False) -> None:
	"""Apply the repository font contract to one PowerPoint run."""
	run.text = clean_inline(text)
	run.font.name = URL_FONT_NAME if url_font else FONT_NAME
	run.font.size = Pt(css_px_to_pt(size))
	run.font.bold = bold
	run.font.color.rgb = color


#============================================
def add_inline_runs(paragraph: object, text: str, size: float, color: object = FOREGROUND,
		url_list: bool = False) -> None:
	"""Write editable bold text, links, and displayed URLs."""
	pattern = re.compile(r"\*\*(.+?)\*\*|\[([^]]+)\]\((https?://[^)]+)\)|(https?://[^\s)]+)")
	position = 0
	for match in pattern.finditer(text):
		if match.start() > position:
			write_run(paragraph.add_run(), text[position:match.start()], size, color)
		if match.group(1) is not None:
			write_run(paragraph.add_run(), match.group(1), size, color, bold=True)
		else:
			label = match.group(2) or match.group(4)
			address = match.group(3) or match.group(4)
			run = paragraph.add_run()
			write_run(run, label, size, ACCENT, url_font=url_list or address == label)
			run.hyperlink.address = address
			run.font.underline = True
		position = match.end()
	if position < len(text):
		write_run(paragraph.add_run(), text[position:], size, color)
	if not paragraph.runs:
		write_run(paragraph.add_run(), text, size, color)


#============================================
def parse_blocks(content: str) -> tuple[
	list[tuple[str, str]], list[tuple[str, int, bool, bool]], list[tuple[str, str]],
]:
	"""Parse supported headings, lists, paragraphs, and images."""
	headings: list[tuple[str, str]] = []
	items: list[tuple[str, int, bool, bool]] = []
	images: list[tuple[str, str]] = []
	for line in content.splitlines():
		stripped = line.strip()
		if not stripped or stripped.startswith("Continuation:"):
			continue
		for image_match in IMAGE_PATTERN.finditer(stripped):
			images.append((image_match.group(1), image_match.group(2)))
		stripped = IMAGE_PATTERN.sub("", stripped).strip()
		if not stripped or stripped.startswith(">"):
			continue
		heading_match = HEADING_PATTERN.match(stripped)
		if heading_match:
			headings.append((heading_match.group(1), heading_match.group(2)))
			continue
		list_match = LIST_PATTERN.match(line)
		if list_match:
			indent = len(list_match.group(1).expandtabs(2)) // 2
			items.append((list_match.group(3), indent, list_match.group(2)[0].isdigit(), False))
		else:
			items.append((stripped, 0, False, True))
	return headings, items, images


#============================================
def wrapped_line_count(text: str, size: float, width: float, level: int = 0) -> int:
	"""Estimate wrapped editable lines with a conservative native-text width."""
	plain_text = clean_inline(text)
	available_width = max(width - level * size * 1.43, size * 3)
	characters_per_line = max(int(available_width / (size * 0.54)), 1)
	words = plain_text.split()
	if not words:
		return 1
	line_count = 1
	current_width = 0
	for word in words:
		word_width = len(word)
		if current_width and current_width + 1 + word_width > characters_per_line:
			line_count += 1
			current_width = word_width
		else:
			current_width += word_width + (1 if current_width else 0)
		if word_width > characters_per_line:
			line_count += (word_width - 1) // characters_per_line
	return line_count


#============================================
def estimate_items_height(items: list[tuple[str, int, bool, bool]], size: float,
		width: float) -> float:
	"""Estimate the CSS-pixel height for native editable list paragraphs."""
	height = 0.0
	for text, level, _, _ in items:
		line_count = wrapped_line_count(text, size, width, level)
		height += line_count * size * BODY_LINE_HEIGHT
		height += size * LIST_ITEM_SPACE_EM
	return height


#============================================
def fit_body_size(item_sets: list[list[tuple[str, int, bool, bool]]], width: float,
		height: float, preferred_size: float, context: str) -> float:
	"""Return one readable CSS size that fits every peer editable text frame."""
	for quarter_points in range(int(preferred_size * 4), int(MIN_READABLE_BODY_SIZE * 4) - 1, -1):
		size = quarter_points / 4
		if all(estimate_items_height(items, size, width) <= height for items in item_sets):
			return size
	raise ValueError(
		f"{context} content cannot fit within the supported readable minimum of "
		f"{MIN_READABLE_BODY_SIZE:g} CSS px",
	)


#============================================
def list_item_groups(items: list[tuple[str, int, bool, bool]]) -> list[list[tuple[str, int, bool, bool]]]:
	"""Attach nested list items to their preceding top-level parent item."""
	groups: list[list[tuple[str, int, bool, bool]]] = []
	for item in items:
		if item[1] == 0 or not groups:
			groups.append([item])
		else:
			groups[-1].append(item)
	return groups


#============================================
def balance_list_columns(items: list[tuple[str, int, bool, bool]], size: float,
		width: float) -> tuple[list[tuple[str, int, bool, bool]], list[tuple[str, int, bool, bool]]]:
	"""Split ordered top-level item groups into two columns of similar height."""
	groups = list_item_groups(items)
	if len(groups) < 2:
		return items, []
	group_heights = [estimate_items_height(group, size, width) for group in groups]
	best_boundary = 1
	best_difference = float("inf")
	left_height = 0.0
	total_height = sum(group_heights)
	for boundary in range(1, len(groups)):
		left_height += group_heights[boundary - 1]
		difference = abs(left_height - (total_height - left_height))
		if difference < best_difference:
			best_boundary = boundary
			best_difference = difference
	left_items = [item for group in groups[:best_boundary] for item in group]
	right_items = [item for group in groups[best_boundary:] for item in group]
	return left_items, right_items


#============================================
def write_items(text_frame: object, items: list[tuple[str, int, bool, bool]], size: float,
		url_list: bool, first_paragraph: object | None = None) -> None:
	"""Write paragraphs and native bullets or continuing automatic numbers."""
	for index, (text, level, ordered, paragraph_only) in enumerate(items):
		paragraph = first_paragraph if index == 0 and first_paragraph is not None else (
			text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
		)
		paragraph.level = level
		paragraph.space_after = Pt(css_px_to_pt(size * LIST_ITEM_SPACE_EM))
		paragraph.line_spacing = BODY_LINE_HEIGHT
		if paragraph_only:
			bullet = OxmlElement("a:buNone")
		elif ordered:
			bullet = OxmlElement("a:buAutoNum")
			bullet.set("type", "arabicPeriod")
		else:
			bullet = OxmlElement("a:buChar")
			bullet.set("char", "\u2022")
		paragraph._p.get_or_add_pPr().insert(0, bullet)
		add_inline_runs(paragraph, text, size, url_list=url_list)


#============================================
def add_title(slide: object, title: str, size: float = 48, url_list: bool = False,
		margin_after: float = 24) -> float:
	"""Add a wrapping native title and return the CSS-pixel content start."""
	title_height = wrapped_line_count(title, size, 1160) * size * 1.12
	text_frame = add_textbox(slide, 60, 52, 1160, title_height)
	paragraph = text_frame.paragraphs[0]
	add_inline_runs(paragraph, title, size, url_list=url_list)
	for run in paragraph.runs:
		run.font.bold = True
	content_top = 52 + title_height + margin_after
	return content_top


#============================================
def parse_panes(content: str) -> list[str]:
	"""Return the two top-level quote blocks used by a two-pane slide."""
	panes = re.split(r"\n\s*\n(?=>)", content.strip())
	result: list[str] = []
	for pane in panes:
		if pane.lstrip().startswith(">"):
			pane_lines: list[str] = []
			for line in pane.splitlines():
				quote_line = line.lstrip()
				content_line = quote_line[1:]
				if content_line.startswith(" "):
					content_line = content_line[1:]
				pane_lines.append(content_line)
			result.append("\n".join(pane_lines).strip())
	return result


#============================================
def add_background(slide: object) -> None:
	"""Set a white canvas and native accent rule."""
	slide.background.fill.solid()
	slide.background.fill.fore_color.rgb = WHITE
	accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(0), px(0), px(1280), px(10))
	accent.fill.solid()
	accent.fill.fore_color.rgb = ACCENT
	accent.line.fill.background()


#============================================
def split_images(images: list[tuple[str, str]]) -> tuple[
		list[tuple[str, str, str, float]], list[tuple[str, str]],
]:
	"""Separate validated bg left/right components from ordinary components."""
	background_images: list[tuple[str, str, str, float]] = []
	ordinary_images: list[tuple[str, str]] = []
	for alt_text, image_value in images:
		if alt_text.startswith("bg"):
			match = BG_IMAGE_PATTERN.fullmatch(alt_text)
			if match is None:
				raise ValueError("bg images require `bg left|right:PERCENT contain` syntax")
			percentage = float(match.group(2))
			if not 1 <= percentage < 100:
				raise ValueError("bg image percentage must be between 1 and 99")
			background_images.append((alt_text, image_value, match.group(1), percentage))
		else:
			ordinary_images.append((alt_text, image_value))
	return background_images, ordinary_images


#============================================
def render_slide(slide: object, source: Slide, deck: Deck) -> None:
	"""Render one supported source slide as editable Office objects."""
	add_background(slide)
	headings, items, images = parse_blocks(source.content)
	title = headings[0][1] if headings and headings[0][0] == "#" else ""
	if "lead" in source.classes:
		text_frame = add_textbox(slide, 110, 180, 1060, 390, MSO_ANCHOR.MIDDLE)
		for index, (_, text) in enumerate(headings):
			paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
			paragraph.alignment = PP_ALIGN.CENTER
			paragraph.space_after = Pt(14)
			add_inline_runs(paragraph, text, 60 if index == 0 else 31,
				FOREGROUND if index == 0 else MUTED, "url-list" in source.classes)
			for run in paragraph.runs:
				run.font.bold = index == 0
		return
	content_top = 128.0
	content_bottom = 754.0
	if title:
		title_margin = 26 if "two-pane" in source.classes else 8 if (
			"figure" in source.classes or "gallery" in source.classes
		) else 24
		content_top = add_title(slide, title, 42 if "dense" in source.classes else 48,
			"url-list" in source.classes, title_margin)
	content_height = content_bottom - content_top
	if "two-pane" in source.classes:
		panes = parse_panes(source.content)
		if len(panes) != 2:
			raise ValueError("two-pane slides require exactly two blockquote panes")
		parsed_panes = [parse_blocks(pane) for pane in panes]
		body_size = 18 if "dense" in source.classes else 22
		pane_heading_size = 28 if "dense" in source.classes else 34
		pane_body_height = content_height
		pane_item_sets = [pane_items for _, pane_items, _ in parsed_panes]
		if any(pane_headings for pane_headings, _, _ in parsed_panes):
			pane_body_height -= pane_heading_size * 1.2 + 10
		body_size = fit_body_size(pane_item_sets, 539, pane_body_height, body_size,
			"two-pane")
		for index, (pane_headings, pane_items, pane_images) in enumerate(parsed_panes):
			left = 60 + index * 581
			text_frame = add_textbox(slide, left, content_top, 539, content_height)
			if pane_headings:
				pane_heading = text_frame.paragraphs[0]
				add_inline_runs(pane_heading, pane_headings[0][1],
					pane_heading_size, FOREGROUND,
					"url-list" in source.classes)
				for run in pane_heading.runs:
					run.font.bold = True
				pane_heading.space_after = Pt(css_px_to_pt(10))
				pane_heading.line_spacing = 1.2
				if pane_items:
					write_items(text_frame, pane_items, body_size,
						"url-list" in source.classes,
						text_frame.add_paragraph())
			else:
				write_items(text_frame, pane_items, body_size, "url-list" in source.classes)
			for alt_text, image_value in pane_images:
				add_picture(slide, resolve_image_path(deck, image_value), alt_text, left, content_top,
					539, content_height)
		return
	if "gallery" in source.classes:
		if not images:
			raise ValueError("gallery slides require component images")
		if items:
			raise ValueError("gallery slides support headings and component images only")
		width = (1160 - 18 * (len(images) - 1)) / len(images)
		for index, (alt_text, image_value) in enumerate(images):
			add_picture(slide, resolve_image_path(deck, image_value), alt_text,
				60 + index * (width + 18), content_top, width, content_height)
		return
	background_images, ordinary_images = split_images(images)
	text_width = 1160
	text_left = 60
	if background_images:
		if len(background_images) != 1:
			raise ValueError("a slide supports one bg left/right contain image")
		alt_text, image_value, direction, percentage = background_images[0]
		image_width = 1160 * percentage / 100
		image_left = 60 if direction == "left" else 1220 - image_width
		add_picture(slide, resolve_image_path(deck, image_value), alt_text, image_left, content_top,
			image_width, content_height)
		text_width -= image_width + 34
		if direction == "left":
			text_left += image_width + 34
	if "figure" in source.classes:
		if not ordinary_images and not background_images:
			raise ValueError("figure slides require a component image")
		if items:
			raise ValueError("figure slides support headings and component images only")
		for alt_text, image_value in ordinary_images:
			add_picture(slide, resolve_image_path(deck, image_value), alt_text, 60, content_top,
				text_width, content_height)
		return
	if ordinary_images and (items or background_images or len(ordinary_images) != 1):
		raise ValueError(
			"ordinary slides support one component image without body text; "
			"use figure, gallery, or bg left/right contain for a multi-content layout",
		)
	size = 20 if "dense" in source.classes or "list-columns" in source.classes else 26
	if "list-columns" in source.classes:
		column_width = text_width / 2 - 21
		columns = balance_list_columns(items, size, column_width)
		size = fit_body_size(list(columns), column_width, content_height, size, "list-columns")
		for index, column_items in enumerate(columns):
			left = text_left + index * (text_width / 2 + 21)
			text_frame = add_textbox(slide, left, content_top, column_width,
				content_height)
			write_items(text_frame, column_items, size, "url-list" in source.classes)
	else:
		text_frame = add_textbox(slide, text_left, content_top, text_width, content_height)
		write_items(text_frame, items, size, "url-list" in source.classes)
	for alt_text, image_value in ordinary_images:
		add_picture(slide, resolve_image_path(deck, image_value), alt_text, text_left, content_top,
			text_width, content_height)


#============================================
def render_native_pptx(deck: Deck, output_path: pathlib.Path) -> pathlib.Path:
	"""Write every canonical slide as native editable PPTX objects."""
	presentation = Presentation()
	presentation.slide_width = px(SLIDE_WIDTH)
	presentation.slide_height = px(SLIDE_HEIGHT)
	presentation.core_properties.title = deck.title
	blank_layout = presentation.slide_layouts[6]
	for number, source in enumerate(deck.slides, start=1):
		slide = presentation.slides.add_slide(blank_layout)
		render_slide(slide, source, deck)
		if source.paginate:
			text_frame = add_textbox(slide, 1190, 762, 62, 22)
			text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
			write_run(text_frame.paragraphs[0].add_run(), str(number), 18, MUTED)
		if source.notes:
			slide.notes_slide.notes_text_frame.text = source.notes
	output_path.parent.mkdir(parents=True, exist_ok=True)
	presentation.save(output_path)
	return output_path


#============================================
def require_executable(command: str, install_message: str) -> pathlib.Path:
	"""Resolve a required executable from PATH."""
	command_value = shutil.which(command)
	if command_value is None:
		raise RuntimeError(install_message)
	command_path = pathlib.Path(command_value).resolve()
	return command_path


#============================================
def convert_pptx(pptx_path: pathlib.Path, output_path: pathlib.Path, output_format: str,
		repo_root: pathlib.Path) -> None:
	"""Use LibreOffice to create an editable ODP or PDF."""
	soffice_path = require_executable("soffice", "LibreOffice is not installed; run brew bundle")
	output_root = repo_root / "output"
	output_root.mkdir(parents=True, exist_ok=True)
	output_path.parent.mkdir(parents=True, exist_ok=True)
	with tempfile.TemporaryDirectory(prefix=".libreoffice.", dir=output_root) as temporary_value:
		temporary_root = pathlib.Path(temporary_value)
		profile_path = temporary_root / "profile"
		conversion_path = temporary_root / "converted"
		conversion_path.mkdir()
		command = [
			str(soffice_path), f"-env:UserInstallation={profile_path.as_uri()}", "--headless",
			"--convert-to", output_format, "--outdir", str(conversion_path), str(pptx_path),
		]
		subprocess.run(command, check=True)
		converted_path = conversion_path / f"{pptx_path.stem}.{output_format}"
		if not converted_path.is_file():
			raise RuntimeError(f"LibreOffice did not create the expected {output_format.upper()} output")
		os.replace(converted_path, output_path)


#============================================
def export_deck(input_value: str, output_format: str) -> dict[str, pathlib.Path]:
	"""Export one authoritative Marp deck without a browser dependency."""
	if output_format not in ("all", "odp", "pdf", "pptx"):
		raise ValueError(f"unsupported output format: {output_format}")
	repo_root = find_repo_root()
	input_path = validate_input(input_value, repo_root)
	deck_name = input_path.stem
	outputs = {
		"pdf": repo_root / f"output/pdf/{deck_name}.pdf",
		"pptx": repo_root / f"output/pptx/{deck_name}.pptx",
		"odp": repo_root / f"output/odp/{deck_name}.odp",
	}
	generated: dict[str, pathlib.Path] = {}
	render_native_pptx(parse_deck(input_path), outputs["pptx"])
	generated["pptx"] = outputs["pptx"]
	if output_format in ("all", "odp"):
		convert_pptx(outputs["pptx"], outputs["odp"], "odp", repo_root)
		generated["odp"] = outputs["odp"]
	if output_format in ("all", "pdf"):
		convert_pptx(outputs["pptx"], outputs["pdf"], "pdf", repo_root)
		generated["pdf"] = outputs["pdf"]
	return generated


#============================================
def print_outputs(outputs: dict[str, pathlib.Path]) -> None:
	"""Print generated artifact paths in classroom workflow order."""
	for output_format in ("pdf", "pptx", "odp"):
		if output_format in outputs:
			print(f"{output_format.upper()}: {outputs[output_format]}")


#============================================
