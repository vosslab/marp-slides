#!/usr/bin/env python3
"""Convert a trusted PPTX into simple, editable Marp Markdown layouts."""

# Standard Library
import io
import os
import re
import json
import zipfile
import pathlib
import argparse
import hashlib
import tempfile
import dataclasses

# PIP3 modules
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


MAX_INPUT_BYTES = 256 * 1024 * 1024
MAX_MEMBER_BYTES = 128 * 1024 * 1024
MAX_UNPACKED_BYTES = 512 * 1024 * 1024
MAX_ARCHIVE_MEMBERS = 4000
MAX_IMAGE_PIXELS = 100_000_000
SUPPORTED_IMAGE_SUFFIXES = {".gif", ".jpeg", ".jpg", ".png"}
SAFE_LINK_SCHEMES = ("http://", "https://", "mailto:")


@dataclasses.dataclass(frozen=True)
class TextBlock:
	"""One positioned PPTX text shape."""

	left: int
	top: int
	lines: tuple[tuple[int, str], ...]
	is_subtitle: bool


@dataclasses.dataclass(frozen=True)
class ImageAsset:
	"""One extracted content image and its source geometry."""

	left: int
	top: int
	width: int
	height: int
	markdown_path: str
	alt_text: str


@dataclasses.dataclass(frozen=True)
class SlideData:
	"""Semantic content extracted from one PPTX slide."""

	source_index: int
	hidden: bool
	title_lines: tuple[str, ...]
	text_blocks: tuple[TextBlock, ...]
	images: tuple[ImageAsset, ...]
	notes: tuple[str, ...]
	review_reasons: tuple[str, ...]


@dataclasses.dataclass(frozen=True)
class ConversionSummary:
	"""User-facing counts for one conversion."""

	visible_slides: int
	editable_slides: int
	hidden_slides: int
	extracted_images: int
	review_slides: int
	output_path: pathlib.Path
	report_path: pathlib.Path


#============================================
def validate_member_name(member_name: str) -> None:
	"""Reject absolute and traversal paths in an OOXML archive."""
	# ASVS 5.3.3: archive member paths never control filesystem destinations.
	normalized = member_name.replace("\\", "/")
	parts = pathlib.PurePosixPath(normalized).parts
	if normalized.startswith("/") or ".." in parts:
		raise ValueError(f"unsafe archive member path: {member_name}")


#============================================
def validate_pptx(input_path: pathlib.Path) -> None:
	"""Validate a bounded OOXML presentation before parsing it."""
	# ASVS 2.2.1 and 5.2.1: enforce the expected type and processing limits.
	if not input_path.is_file() or input_path.suffix.lower() != ".pptx":
		raise ValueError("input must be an existing .pptx file")
	if input_path.stat().st_size > MAX_INPUT_BYTES:
		raise ValueError("PPTX exceeds the compressed input limit")
	with zipfile.ZipFile(input_path) as archive:
		members = archive.infolist()
		if len(members) > MAX_ARCHIVE_MEMBERS:
			raise ValueError("PPTX contains too many archive members")
		total_size = 0
		member_names: set[str] = set()
		for member in members:
			validate_member_name(member.filename)
			if member.file_size > MAX_MEMBER_BYTES:
				raise ValueError(f"PPTX member exceeds size limit: {member.filename}")
			total_size += member.file_size
			if total_size > MAX_UNPACKED_BYTES:
				raise ValueError("PPTX exceeds the expanded archive limit")
			member_names.add(member.filename)
		required_names = {"[Content_Types].xml", "ppt/presentation.xml"}
		if not required_names.issubset(member_names):
			raise ValueError("PPTX is missing required OOXML presentation members")


#============================================
def markdown_text(value: str) -> str:
	"""Encode source text only at the final Markdown output boundary."""
	# ASVS 1.1.2 and 1.2.1: escape for the destination syntax at render time.
	encoded = value.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
	for character in ("\\", "`", "*", "[", "]"):
		encoded = encoded.replace(character, f"\\{character}")
	return encoded


#============================================
def comment_text(value: str) -> str:
	"""Encode one line for a Marp presenter-note HTML comment."""
	encoded = value.replace("--", "- -")
	encoded = encoded.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
	encoded = encoded.encode("ascii", "xmlcharrefreplace").decode("ascii")
	return encoded


#============================================
def safe_hyperlink(address: str | None) -> str:
	"""Return an allowed hyperlink address or an empty string."""
	if address is None:
		return ""
	# ASVS 1.2.2: only known-safe URL schemes enter generated Markdown links.
	if not address.lower().startswith(SAFE_LINK_SCHEMES):
		return ""
	return address.replace(" ", "%20").replace("(", "%28").replace(")", "%29")


#============================================
def paragraph_markdown(paragraph: object) -> str:
	"""Render one python-pptx paragraph without losing its hyperlinks."""
	segments: list[str] = []
	for run in paragraph.runs:
		text = markdown_text(run.text)
		if not text:
			continue
		if segments and segments[-1][-1:].isalnum() and text[:1].isalnum():
			segments.append(" ")
		address = safe_hyperlink(run.hyperlink.address)
		if address:
			segments.append(f"[{text}]({address})")
		else:
			segments.append(text)
	if segments:
		return "".join(segments).strip()
	return markdown_text(paragraph.text).strip()


#============================================
def title_shape_id(slide: object) -> int | None:
	"""Return the title shape id without relying on proxy identity."""
	shape = slide.shapes.title
	if shape is None:
		return None
	return shape.shape_id


#============================================
def is_subtitle_shape(shape: object) -> bool:
	"""Return whether a placeholder is a subtitle."""
	if not shape.is_placeholder:
		return False
	placeholder_name = str(shape.placeholder_format.type).upper()
	return "SUBTITLE" in placeholder_name


#============================================
def validate_image_blob(blob: bytes, suffix: str) -> None:
	"""Validate one extracted raster image before writing it."""
	# ASVS 5.2.2 and 5.2.6: validate the declared type and decoded pixel bounds.
	if suffix not in SUPPORTED_IMAGE_SUFFIXES:
		raise ValueError(f"unsupported PPTX image type: {suffix}")
	if len(blob) > MAX_MEMBER_BYTES:
		raise ValueError("PPTX image exceeds the per-image size limit")
	with Image.open(io.BytesIO(blob)) as image:
		width, height = image.size
		if width < 1 or height < 1 or width * height > MAX_IMAGE_PIXELS:
			raise ValueError("PPTX image dimensions exceed the supported limit")
		image.verify()


#============================================
def image_asset(
	shape: object,
	assets_dir: pathlib.Path,
	markdown_root: pathlib.PurePosixPath,
	known_images: dict[str, str],
) -> ImageAsset:
	"""Extract one picture under an internally generated content hash."""
	blob = shape.image.blob
	suffix = f".{shape.image.ext.lower()}"
	if suffix == ".jpeg":
		suffix = ".jpg"
	validate_image_blob(blob, suffix)
	digest = hashlib.sha256(blob).hexdigest()
	asset_name = known_images.get(digest)
	if asset_name is None:
		asset_name = f"image_{len(known_images) + 1:03d}{suffix}"
		# ASVS 5.3.2: source filenames never control the destination path.
		(assets_dir / asset_name).write_bytes(blob)
	known_images[digest] = asset_name
	markdown_path = (markdown_root / asset_name).as_posix()
	asset_number = int(re.search(r"\d+", asset_name).group())
	return ImageAsset(
		left=shape.left,
		top=shape.top,
		width=shape.width,
		height=shape.height,
		markdown_path=markdown_path,
		alt_text=f"Slide image {asset_number}",
	)


#============================================
def shape_inventory(
	shape: object,
	title_id: int | None,
	assets_dir: pathlib.Path,
	markdown_root: pathlib.PurePosixPath,
	known_images: dict[str, str],
) -> tuple[list[str], list[TextBlock], list[ImageAsset], list[str]]:
	"""Extract supported semantic objects recursively from one shape."""
	titles: list[str] = []
	text_blocks: list[TextBlock] = []
	images: list[ImageAsset] = []
	review_reasons: list[str] = []
	if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
		for child in shape.shapes:
			child_parts = shape_inventory(
				child,
				title_id,
				assets_dir,
				markdown_root,
				known_images,
			)
			titles.extend(child_parts[0])
			text_blocks.extend(child_parts[1])
			images.extend(child_parts[2])
			review_reasons.extend(child_parts[3])
		return titles, text_blocks, images, review_reasons
	if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
		images.append(image_asset(shape, assets_dir, markdown_root, known_images))
		return titles, text_blocks, images, review_reasons
	if getattr(shape, "has_table", False):
		for row in shape.table.rows:
			cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
			if cells:
				line = " | ".join(markdown_text(cell) for cell in cells)
				text_blocks.append(TextBlock(shape.left, shape.top, ((0, line),), False))
		return titles, text_blocks, images, review_reasons
	if getattr(shape, "has_text_frame", False):
		line_values: list[tuple[int, str]] = []
		for paragraph in shape.text_frame.paragraphs:
			text = paragraph_markdown(paragraph)
			if text:
				line_values.append((paragraph.level, text))
		lines = tuple(line_values)
		if not lines:
			if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
				review_reasons.append(f"ignored non-content shape type {shape.shape_type}")
			return titles, text_blocks, images, review_reasons
		if shape.shape_id == title_id:
			titles.extend(text for _level, text in lines)
		else:
			text_blocks.append(
				TextBlock(shape.left, shape.top, lines, is_subtitle_shape(shape))
			)
		return titles, text_blocks, images, review_reasons
	if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
		review_reasons.append(f"ignored non-content shape type {shape.shape_type}")
	return titles, text_blocks, images, review_reasons


#============================================
def slide_notes(slide: object) -> tuple[str, ...]:
	"""Read presenter notes only when the source slide owns a notes part."""
	has_notes_part = any(rel.reltype.endswith("/notesSlide") for rel in slide.part.rels.values())
	if not has_notes_part:
		return ()
	note_text = slide.notes_slide.notes_text_frame.text
	lines = tuple(line.strip() for line in note_text.splitlines() if line.strip())
	return lines


#============================================
def extract_slides(
	presentation: object,
	assets_dir: pathlib.Path,
	markdown_root: pathlib.PurePosixPath,
	expected_hidden: set[int] | None,
) -> tuple[list[SlideData], int]:
	"""Extract all slides while preserving source order and visibility."""
	known_images: dict[str, str] = {}
	slides: list[SlideData] = []
	for source_index, slide in enumerate(presentation.slides, start=1):
		pptx_hidden = slide.element.get("show") == "0"
		hidden = pptx_hidden if expected_hidden is None else source_index in expected_hidden
		if expected_hidden is not None and pptx_hidden != hidden:
			raise RuntimeError(f"ODP and PPTX visibility disagree on slide {source_index}")
		if hidden:
			slides.append(SlideData(source_index, True, (), (), (), (), ()))
			continue
		titles: list[str] = []
		text_blocks: list[TextBlock] = []
		images: list[ImageAsset] = []
		review_reasons: list[str] = []
		title_id = title_shape_id(slide)
		for shape in slide.shapes:
			parts = shape_inventory(
				shape,
				title_id,
				assets_dir,
				markdown_root,
				known_images,
			)
			titles.extend(parts[0])
			text_blocks.extend(parts[1])
			images.extend(parts[2])
			review_reasons.extend(parts[3])
		line_count = sum(len(block.lines) for block in text_blocks)
		character_count = sum(len(text) for block in text_blocks for _level, text in block.lines)
		if line_count > 10 or character_count > 1200:
			review_reasons.append("dense text requires post-conversion polish")
		text_blocks.sort(key=lambda block: (block.top, block.left))
		images.sort(key=lambda image: (image.top, image.left))
		slides.append(
			SlideData(
				source_index=source_index,
				hidden=hidden,
				title_lines=tuple(titles),
				text_blocks=tuple(text_blocks),
				images=tuple(images),
				notes=slide_notes(slide),
				review_reasons=tuple(sorted(set(review_reasons))),
			)
		)
	return slides, len(known_images)


#============================================
def body_lines(slide: SlideData) -> list[str]:
	"""Render positioned text blocks as a simple semantic outline."""
	lines: list[str] = []
	for block in slide.text_blocks:
		for line_index, (level, text) in enumerate(block.lines):
			if block.is_subtitle:
				lines.append(f"## {text}")
			else:
				lines.append(f"{'  ' * level}- {text}")
	return lines


#============================================
def quoted(lines: list[str]) -> list[str]:
	"""Wrap ordinary Markdown lines in one blockquote pane."""
	return [f"> {line}" if line else ">" for line in lines]


#============================================
def image_markdown(image: ImageAsset) -> str:
	"""Render one ordinary content image."""
	return f"![{image.alt_text}]({image.markdown_path})"


#============================================
def render_slide(slide: SlideData, slide_width: int, is_first: bool) -> tuple[list[str], str]:
	"""Map one semantic slide into the smallest suitable Marp layout."""
	texts = body_lines(slide)
	title_lines = list(slide.title_lines)
	if not title_lines and texts:
		first_text = texts.pop(0).removeprefix("- ")
		title_lines = [first_text]
	if not title_lines:
		title_lines = [f"Slide {slide.source_index}"]
	heading_lines = [f"# {' '.join(title_lines)}"]
	images = list(slide.images)
	if not images:
		lead = is_first or not texts
		prefix = ["<!-- _class: lead -->", "<!-- _paginate: false -->"] if lead else []
		return [*prefix, *heading_lines, "", *texts], "lead" if lead else "default"
	if len(images) == 1 and not texts:
		return ["<!-- _class: figure -->", *heading_lines, "", image_markdown(images[0])], "figure"
	if len(images) == 1:
		image = images[0]
		center = (image.left + image.width / 2) / slide_width
		pane_percent = max(28, min(55, round(100 * image.width / slide_width)))
		if center >= 0.55:
			directive = f"![bg right:{pane_percent}% contain]({image.markdown_path})"
			return [*heading_lines, "", *texts, "", directive], "right-image"
		if center <= 0.45:
			directive = f"![bg left:{pane_percent}% contain]({image.markdown_path})"
			return [*heading_lines, "", *texts, "", directive], "left-image"
	image_lines = [image_markdown(image) for image in images]
	if not texts:
		return [
			"<!-- _class: gallery -->",
			*heading_lines,
			"",
			" ".join(image_lines),
		], "gallery"
	return [
		"<!-- _class: two-pane -->",
		*heading_lines,
		"",
		*quoted(texts),
		"",
		*quoted([" ".join(image_lines)]),
	], "two-pane"


#============================================
def render_markdown(
	slides: list[SlideData],
	slide_width: int,
	title: str,
) -> tuple[str, list[dict[str, object]]]:
	"""Render one complete Marp deck and its per-slide report records."""
	frontmatter = [
		"---",
		"marp: true",
		"theme: genetics",
		"size: 16:10",
		"paginate: true",
		f"title: {json.dumps(title, ensure_ascii=True)}",
		"---",
	]
	visible_slides = [slide for slide in slides if not slide.hidden]
	markdown_lines = frontmatter
	records: list[dict[str, object]] = []
	for visible_index, slide in enumerate(visible_slides):
		markdown_lines.extend(["", "---", ""] if visible_index else [""])
		slide_lines, layout = render_slide(slide, slide_width, visible_index == 0)
		markdown_lines.extend(slide_lines)
		if slide.notes:
			markdown_lines.extend(["", "<!--"])
			markdown_lines.extend(comment_text(line) for line in slide.notes)
			markdown_lines.append("-->")
		records.append(
			{
				"source_slide": slide.source_index,
				"layout": layout,
				"text_blocks": len(slide.text_blocks),
				"images": len(slide.images),
				"review_reasons": list(slide.review_reasons),
			}
		)
	return "\n".join(markdown_lines).rstrip() + "\n", records


#============================================
def validate_output_path(output_path: pathlib.Path) -> None:
	"""Protect an established canonical destination from replacement."""
	if output_path.suffix.lower() != ".md":
		raise ValueError("output must use the .md extension")
	if output_path.exists():
		raise FileExistsError("output Markdown already exists; import will not overwrite it")
	asset_path = output_path.parent / "assets" / output_path.stem
	if asset_path.exists():
		raise FileExistsError("output asset directory already exists")


#============================================
def convert_pptx(
	input_path: pathlib.Path,
	output_path: pathlib.Path,
	*,
	expected_slide_count: int | None = None,
	expected_hidden: set[int] | None = None,
	source_name: str | None = None,
) -> ConversionSummary:
	"""Convert one trusted PPTX into a new Marp deck and content assets."""
	input_path = input_path.resolve()
	output_path = output_path.resolve()
	validate_pptx(input_path)
	validate_output_path(output_path)
	presentation = Presentation(input_path)
	if expected_slide_count is not None and len(presentation.slides) != expected_slide_count:
		raise RuntimeError("ODP and normalized PPTX slide counts disagree")
	output_path.parent.mkdir(parents=True, exist_ok=True)
	with tempfile.TemporaryDirectory(
		prefix=".pptx_to_marp_",
		dir=output_path.parent,
	) as temporary_name:
		temporary_root = pathlib.Path(temporary_name)
		staging_assets = temporary_root / "assets"
		staging_assets.mkdir()
		markdown_root = pathlib.PurePosixPath("assets") / output_path.stem
		slides, image_count = extract_slides(
			presentation,
			staging_assets,
			markdown_root,
			expected_hidden,
		)
		visible_slides = [slide for slide in slides if not slide.hidden]
		if not visible_slides:
			raise ValueError("presentation contains no visible slides")
		if visible_slides[0].title_lines:
			deck_title = visible_slides[0].title_lines[0]
		else:
			deck_title = output_path.stem
		markdown, records = render_markdown(slides, presentation.slide_width, deck_title)
		report = {
			"source": source_name or input_path.name,
			"slide_count": len(slides),
			"visible_slides": len(visible_slides),
			"hidden_slides": [slide.source_index for slide in slides if slide.hidden],
			"slides": records,
		}
		report_path = staging_assets / "import_report.json"
		report_path.write_text(json.dumps(report, indent=2) + "\n", encoding="utf-8")
		staging_markdown = temporary_root / output_path.name
		staging_markdown.write_text(markdown, encoding="utf-8")
		final_assets = output_path.parent / "assets" / output_path.stem
		final_assets.parent.mkdir(exist_ok=True)
		os.replace(staging_assets, final_assets)
		os.replace(staging_markdown, output_path)
	final_report = output_path.parent / "assets" / output_path.stem / "import_report.json"
	review_count = sum(bool(slide.review_reasons) for slide in slides if not slide.hidden)
	return ConversionSummary(
		visible_slides=len(visible_slides),
		editable_slides=len(visible_slides),
		hidden_slides=len(slides) - len(visible_slides),
		extracted_images=image_count,
		review_slides=review_count,
		output_path=output_path,
		report_path=final_report,
	)


#============================================
def parse_args() -> argparse.Namespace:
	"""Parse the standalone PPTX importer arguments."""
	parser = argparse.ArgumentParser(description=__doc__)
	parser.add_argument("input_file", type=pathlib.Path, help="trusted source PPTX")
	parser.add_argument("--output", type=pathlib.Path, help="new Marp Markdown path")
	return parser.parse_args()


#============================================
def main() -> None:
	"""Run one standalone PPTX-to-Marp conversion."""
	args = parse_args()
	output_path = args.output or args.input_file.with_suffix(".md")
	summary = convert_pptx(args.input_file, output_path)
	print(
		f"Converted {summary.visible_slides} visible slides: "
		f"{summary.editable_slides} editable, {summary.review_slides} layout review, "
		f"{summary.hidden_slides} hidden, {summary.extracted_images} content images"
	)
	print(f"Markdown: {summary.output_path}")
	print(f"Import report: {summary.report_path}")


if __name__ == "__main__":
	main()
